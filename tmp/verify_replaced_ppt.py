from pathlib import Path

from PIL import Image, ImageDraw, ImageFont
from pptx import Presentation


def find_output() -> Path:
    downloads = Path.home() / "Downloads"
    candidates = [
        path
        for path in downloads.glob("*.pptx")
        if "CIM" in path.name and "截图替换" in path.stem and not path.name.startswith("~$")
    ]
    if not candidates:
        raise SystemExit("No replaced PPTX found.")
    return max(candidates, key=lambda path: path.stat().st_mtime)


def picture_shapes(slide):
    for shape in slide.shapes:
        try:
            if shape.shape_type == 13:
                yield shape
        except NotImplementedError:
            continue


out_dir = Path("tmp/pptx_replaced_verify")
out_dir.mkdir(parents=True, exist_ok=True)

target = find_output()
prs = Presentation(target)
print(f"target={target}")
print(f"slides={len(prs.slides)}")

exports = []
for slide_number in (5, 8):
    slide = prs.slides[slide_number - 1]
    for pic_number, shape in enumerate(picture_shapes(slide), 1):
        if pic_number == 1:
            continue
        image = shape.image
        path = out_dir / f"slide{slide_number}_pic{pic_number}.{image.ext or 'png'}"
        path.write_bytes(image.blob)
        exports.append((slide_number, pic_number, path))
        print(f"{path} {image.size}")

cell_w, cell_h = 310, 210
cols = 4
rows = (len(exports) + cols - 1) // cols
sheet = Image.new("RGB", (cols * cell_w, rows * cell_h), "white")
draw = ImageDraw.Draw(sheet)
font = ImageFont.load_default()
for index, (slide_number, pic_number, path) in enumerate(exports):
    img = Image.open(path).convert("RGB")
    img.thumbnail((270, 150))
    col = index % cols
    row = index // cols
    x = col * cell_w
    y = row * cell_h
    draw.rectangle((x, y, x + cell_w - 1, y + cell_h - 1), outline=(205, 205, 205))
    draw.text((x + 10, y + 8), f"slide{slide_number} pic{pic_number}", fill=(0, 0, 0), font=font)
    draw.text((x + 10, y + 24), path.name, fill=(80, 80, 80), font=font)
    sheet.paste(img, (x + (cell_w - img.width) // 2, y + 45 + (150 - img.height) // 2))

sheet_path = out_dir / "contact_sheet.jpg"
sheet.save(sheet_path, quality=92)
print(f"contact_sheet={sheet_path}")
