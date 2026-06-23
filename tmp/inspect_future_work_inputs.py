from __future__ import annotations

from pathlib import Path
import sys

from pptx import Presentation


def find_ppt() -> Path:
    downloads = Path.home() / "Downloads"
    candidates = [
        path
        for path in downloads.glob("*.pptx")
        if "CIM" in path.name and "截图替换" in path.stem and not path.name.startswith("~$")
    ]
    if not candidates:
        raise SystemExit("No replaced PPTX found in Downloads.")
    return max(candidates, key=lambda path: path.stat().st_mtime)


def find_xlsx() -> Path:
    root = Path.home() / "Desktop" / "chk" / "CIMAgent"
    candidates = [path for path in root.glob("*.xlsx") if "曹宏坤" in path.name and not path.name.startswith("~$")]
    if not candidates:
        raise SystemExit("No matching xlsx found.")
    return max(candidates, key=lambda path: path.stat().st_mtime)


def inspect_ppt(path: Path) -> None:
    prs = Presentation(path)
    slide = prs.slides[11]
    print(f"PPT={path}")
    print(f"slides={len(prs.slides)}")
    print("slide12 shapes:")
    for index, shape in enumerate(slide.shapes, 1):
        text = ""
        if getattr(shape, "has_text_frame", False) and shape.has_text_frame:
            text = shape.text.strip().replace("\n", " | ")
        try:
            shape_type = str(shape.shape_type)
        except NotImplementedError:
            shape_type = "unrecognized"
        print(
            f"  {index}: type={shape_type} name={shape.name!r} "
            f"left={shape.left} top={shape.top} width={shape.width} height={shape.height} "
            f"text={text[:240]!r}"
        )


def inspect_xlsx(path: Path) -> None:
    try:
        import openpyxl
    except ImportError as exc:
        print(f"openpyxl missing, using fallback parser: {exc}")
        inspect_xlsx_fallback(path)
        return

    wb = openpyxl.load_workbook(path, data_only=True)
    print(f"XLSX={path}")
    print(f"sheets={wb.sheetnames}")
    for ws in wb.worksheets:
        print(f"sheet={ws.title} max_row={ws.max_row} max_column={ws.max_column}")
        rows_printed = 0
        for row in ws.iter_rows(values_only=True):
            values = ["" if value is None else str(value).strip() for value in row]
            if not any(values):
                continue
            print("  row:", " | ".join(values))
            rows_printed += 1
            if rows_printed >= 80:
                print("  ...")
                break


def inspect_xlsx_fallback(path: Path) -> None:
    from zipfile import ZipFile
    import re
    import xml.etree.ElementTree as ET

    ns = {
        "a": "http://schemas.openxmlformats.org/spreadsheetml/2006/main",
        "r": "http://schemas.openxmlformats.org/officeDocument/2006/relationships",
        "rel": "http://schemas.openxmlformats.org/package/2006/relationships",
    }

    def text_content(element: ET.Element | None) -> str:
        if element is None:
            return ""
        return "".join(element.itertext()).strip()

    def col_index(cell_ref: str) -> int:
        letters = re.match(r"[A-Z]+", cell_ref or "")
        if not letters:
            return 0
        value = 0
        for char in letters.group(0):
            value = value * 26 + ord(char) - ord("A") + 1
        return value - 1

    with ZipFile(path, "r") as zf:
        shared_strings: list[str] = []
        if "xl/sharedStrings.xml" in zf.namelist():
            root = ET.fromstring(zf.read("xl/sharedStrings.xml"))
            shared_strings = [text_content(si) for si in root.findall("a:si", ns)]

        workbook = ET.fromstring(zf.read("xl/workbook.xml"))
        rel_root = ET.fromstring(zf.read("xl/_rels/workbook.xml.rels"))
        rels = {rel.attrib["Id"]: rel.attrib["Target"] for rel in rel_root.findall("rel:Relationship", ns)}
        print(f"XLSX={path}")
        sheets = []
        for sheet in workbook.findall(".//a:sheet", ns):
            name = sheet.attrib.get("name", "")
            rid = sheet.attrib.get(f"{{{ns['r']}}}id", "")
            target = rels.get(rid, "")
            sheet_path = "xl/" + target.lstrip("/")
            sheets.append((name, sheet_path))
        print(f"sheets={[name for name, _ in sheets]}")

        for name, sheet_path in sheets:
            root = ET.fromstring(zf.read(sheet_path))
            rows = []
            for row in root.findall(".//a:sheetData/a:row", ns):
                values: list[str] = []
                for cell in row.findall("a:c", ns):
                    index = col_index(cell.attrib.get("r", ""))
                    while len(values) <= index:
                        values.append("")
                    cell_type = cell.attrib.get("t", "")
                    if cell_type == "s":
                        raw = text_content(cell.find("a:v", ns))
                        value = shared_strings[int(raw)] if raw else ""
                    elif cell_type == "inlineStr":
                        value = text_content(cell.find("a:is", ns))
                    else:
                        value = text_content(cell.find("a:v", ns))
                    values[index] = value
                if any(value.strip() for value in values):
                    rows.append(values)
            max_col = max((len(row) for row in rows), default=0)
            print(f"sheet={name} rows={len(rows)} max_column={max_col}")
            for row in rows[:100]:
                print("  row:", " | ".join(row))
            if len(rows) > 100:
                print("  ...")


def main() -> None:
    sys.stdout.reconfigure(encoding="utf-8")
    inspect_ppt(find_ppt())
    inspect_xlsx(find_xlsx())


if __name__ == "__main__":
    main()
