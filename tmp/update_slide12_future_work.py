from __future__ import annotations

from pathlib import Path
import sys

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Pt


def find_ppt() -> Path:
    downloads = Path.home() / "Downloads"
    candidates = [
        path
        for path in downloads.glob("*.pptx")
        if "CIM" in path.name
        and "截图替换" in path.stem
        and "第12页更新" not in path.stem
        and not path.name.startswith("~$")
    ]
    if not candidates:
        raise SystemExit("No replaced PPTX found in Downloads.")
    return max(candidates, key=lambda path: path.stat().st_mtime)


def first_run(shape):
    if not getattr(shape, "has_text_frame", False) or not shape.has_text_frame:
        return None
    for paragraph in shape.text_frame.paragraphs:
        for run in paragraph.runs:
            return run
    return None


def set_text(shape, text: str, *, size: float | None = None, bold: bool | None = None) -> None:
    old_run = first_run(shape)
    old_font = old_run.font if old_run is not None else None
    color = None
    if old_font is not None:
        try:
            color = old_font.color.rgb
        except AttributeError:
            color = None
    if color is None:
        color = RGBColor(255, 255, 255)

    text_frame = shape.text_frame
    text_frame.clear()
    text_frame.word_wrap = True
    paragraph = text_frame.paragraphs[0]
    paragraph.alignment = PP_ALIGN.CENTER
    run = paragraph.add_run()
    run.text = text
    run.font.name = "Microsoft YaHei"
    run.font.size = Pt(size) if size is not None else (old_font.size if old_font is not None else Pt(14))
    run.font.bold = bold if bold is not None else (old_font.bold if old_font is not None else False)
    run.font.color.rgb = color


def set_body(shape, text: str) -> None:
    text_frame = shape.text_frame
    text_frame.clear()
    text_frame.word_wrap = True
    paragraph = text_frame.paragraphs[0]
    paragraph.alignment = PP_ALIGN.CENTER
    run = paragraph.add_run()
    run.text = text
    run.font.name = "Microsoft YaHei"
    run.font.size = Pt(12)
    run.font.bold = False
    run.font.color.rgb = RGBColor(255, 255, 255)


def main() -> None:
    sys.stdout.reconfigure(encoding="utf-8")
    src = find_ppt()
    prs = Presentation(src)
    slide = prs.slides[11]

    # Existing card text shape indexes from slide 12 inspection.
    title_shape = slide.shapes[2]
    card_titles = {
        "01": slide.shapes[19],
        "02": slide.shapes[26],
        "03": slide.shapes[6],
        "04": slide.shapes[14],
    }
    card_bodies = {
        "01": slide.shapes[20],
        "02": slide.shapes[27],
        "03": slide.shapes[7],
        "04": slide.shapes[15],
    }

    set_text(title_shape, "复盘与下一步：围绕 2026 绩效目标推进", size=18, bold=True)
    updates = {
        "01": (
            "自动化建模交付",
            "完善参数提取、模型生成与规则配置，完成 CIM3/CIM4 轨道交通、地下空间部件建模工具集成测试和课题交付。",
        ),
        "02": (
            "超融合数据库",
            "基于 ClickHouse 完成不少于 6 类多模态时空数据接入、时空索引与组合检索，推进性能优化和软著成果。",
        ),
        "03": (
            "时空基座大模型",
            "围绕 Qwen-VL 开展地理空间与时间约束微调，完成训练样本、评测验证和城市时空基座模型阶段成果。",
        ),
        "04": (
            "成果沉淀与支撑",
            "支撑重点项目、汇报材料和技术验证闭环，同步推进 1 篇 AI 相关 EI 论文、2 项专利及研发资料归档。",
        ),
    }

    for key, (heading, body) in updates.items():
        set_text(card_titles[key], heading, size=15, bold=True)
        set_body(card_bodies[key], body)

    out = src.with_name(f"{src.stem}_第12页更新.pptx")
    prs.save(out)
    print(f"source={src}")
    print(f"output={out}")
    for key, (heading, body) in updates.items():
        print(f"{key} {heading}: {body}")


if __name__ == "__main__":
    main()
