from pathlib import Path

from pptx import Presentation


def find_target() -> Path:
    downloads = Path.home() / "Downloads"
    candidates = [
        path
        for path in downloads.glob("*.pptx")
        if not path.name.startswith("~$")
        and "CIM" in path.name
        and path.name.endswith("11.pptx")
    ]
    if not candidates:
        raise SystemExit("No matching PPTX found in Downloads.")
    return max(candidates, key=lambda path: path.stat().st_mtime)


path = find_target()
prs = Presentation(path)
print(f"target={path}")
print(f"slides={len(prs.slides)}")

for slide_index, slide in enumerate(prs.slides, 1):
    title = ""
    for shape in slide.shapes:
        if getattr(shape, "has_text_frame", False) and shape.has_text_frame:
            text = shape.text.strip()
            if text:
                title = text.replace("\n", " | ")[:100]
                break

    pictures = []
    for shape in slide.shapes:
        try:
            if shape.shape_type == 13:
                pictures.append(shape)
        except NotImplementedError:
            continue
    print(f"Slide {slide_index}: pics={len(pictures)} title={title}")
    if slide_index in (5, 8):
        for pic_index, shape in enumerate(pictures, 1):
            print(
                "  "
                f"pic{pic_index}: name={shape.name!r} "
                f"left={shape.left} top={shape.top} "
                f"width={shape.width} height={shape.height}"
            )
