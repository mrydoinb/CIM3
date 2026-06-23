from pathlib import Path

from PIL import Image, ImageDraw, ImageFont
from pptx import Presentation


def find_target() -> Path:
    downloads = Path.home() / "Downloads"
    candidates = [
        path
        for path in downloads.glob("*.pptx")
        if not path.name.startswith("~$")
        and "CIM" in path.name
        and path.name.endswith("11.pptx")
    ]
    if not candidates:
        raise SystemExit("No matching PPTX found in Downloads.")
    return max(candidates, key=lambda path: path.stat().st_mtime)


def picture_shapes(slide):
    for shape in slide.shapes:
        try:
            if shape.shape_type == 13:
                yield shape
        except NotImplementedError:
            continue


out_dir = Path("tmp/pptx_target_images")
out_dir.mkdir(parents=True, exist_ok=True)

prs = Presentation(find_target())
exports = []
for slide_number in (5, 8):
    slide = prs.slides[slide_number - 1]
    for pic_number, shape in enumerate(picture_shapes(slide), 1):
        image = shape.image
        ext = image.ext or "png"
        out_path = out_dir / f"slide{slide_number}_pic{pic_number}.{ext}"
        out_path.write_bytes(image.blob)
        exports.append((slide_number, pic_number, out_path, shape.width, shape.height))
        print(f"{out_path} {image.size} {shape.width}x{shape.height}")

thumbs = []
for slide_number, pic_number, path, width, height in exports:
    img = Image.open(path).convert("RGB")
    img.thumbnail((260, 160))
    thumbs.append((slide_number, pic_number, img.copy(), path.name))

if thumbs:
    cell_w, cell_h = 310, 210
    cols = 4
    rows = (len(thumbs) + cols - 1) // cols
    sheet = Image.new("RGB", (cols * cell_w, rows * cell_h), "white")
    draw = ImageDraw.Draw(sheet)
    font = ImageFont.load_default()
    for index, (slide_number, pic_number, img, name) in enumerate(thumbs):
        col = index % cols
        row = index // cols
        x = col * cell_w
        y = row * cell_h
        draw.rectangle([x, y, x + cell_w - 1, y + cell_h - 1], outline=(210, 210, 210))
        draw.text((x + 10, y + 8), f"slide{slide_number} pic{pic_number}", fill=(0, 0, 0), font=font)
        draw.text((x + 10, y + 24), name, fill=(80, 80, 80), font=font)
        ix = x + (cell_w - img.width) // 2
        iy = y + 45 + (150 - img.height) // 2
        sheet.paste(img, (ix, iy))
    sheet_path = out_dir / "contact_sheet.jpg"
    sheet.save(sheet_path, quality=92)
    print(f"contact_sheet={sheet_path}")
