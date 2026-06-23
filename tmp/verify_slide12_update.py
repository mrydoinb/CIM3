from __future__ import annotations

from pathlib import Path
import sys

from pptx import Presentation


def find_ppt() -> Path:
    downloads = Path.home() / "Downloads"
    candidates = [
        path
        for path in downloads.glob("*.pptx")
        if "CIM" in path.name and "第12页更新" in path.stem and not path.name.startswith("~$")
    ]
    if not candidates:
        raise SystemExit("No updated PPTX found.")
    return max(candidates, key=lambda path: path.stat().st_mtime)


sys.stdout.reconfigure(encoding="utf-8")
path = find_ppt()
prs = Presentation(path)
print(f"PPT={path}")
print(f"size={path.stat().st_size}")
print(f"slides={len(prs.slides)}")
for index, shape in enumerate(prs.slides[11].shapes, 1):
    if getattr(shape, "has_text_frame", False) and shape.has_text_frame:
        text = shape.text.strip()
        if text:
            print(f"{index}: {shape.name}: {text.replace(chr(10), ' | ')}")
