from io import BytesIO
from pathlib import Path

from PIL import Image, ImageDraw, ImageFont
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_CONNECTOR, MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt


ROOT = Path(r"C:\Users\22838\Desktop\chk\CIMAgent\cim_road_poc")
SOURCE = Path(r"C:\Users\22838\Desktop\chk\CIMAgent\PPT汇报\chkpptx.pptx")
OUTPUT = ROOT / "chkpptx_后续重点工作双专题版.pptx"
PREVIEW_9 = ROOT / "output" / "chkpptx_城市时空基座大模型_第9页预览.png"
PREVIEW_10 = ROOT / "output" / "chkpptx_超融合数据库_第10页预览.png"

FONT = "Microsoft YaHei"

NAVY = RGBColor(16, 50, 88)
DEEP_BLUE = RGBColor(0, 82, 156)
BLUE = RGBColor(16, 112, 224)
CYAN = RGBColor(24, 166, 207)
GREEN = RGBColor(33, 145, 103)
ORANGE = RGBColor(229, 139, 48)
RED = RGBColor(235, 35, 58)
INK = RGBColor(30, 49, 72)
MID = RGBColor(91, 111, 135)
LIGHT = RGBColor(215, 228, 240)
PALE_BLUE = RGBColor(239, 247, 255)
PALE_GREEN = RGBColor(236, 248, 243)
PALE_ORANGE = RGBColor(255, 244, 231)
WHITE = RGBColor(255, 255, 255)


def remove_all_shapes(slide):
    for shape in list(slide.shapes):
        element = shape._element
        element.getparent().remove(element)


def set_background(slide, color=WHITE):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_text(
    slide,
    text,
    x,
    y,
    w,
    h,
    size=12,
    color=INK,
    bold=False,
    align=PP_ALIGN.LEFT,
    valign=MSO_ANCHOR.MIDDLE,
    margin=0.03,
):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    frame = box.text_frame
    frame.clear()
    frame.word_wrap = True
    frame.margin_left = frame.margin_right = Inches(margin)
    frame.margin_top = frame.margin_bottom = Inches(margin)
    frame.vertical_anchor = valign
    paragraph = frame.paragraphs[0]
    paragraph.text = text
    paragraph.alignment = align
    paragraph.space_after = Pt(0)
    paragraph.line_spacing = 1.05
    run = paragraph.runs[0]
    run.font.name = FONT
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    return box


def add_rect(slide, x, y, w, h, fill, line=None, rounded=False):
    shape_type = MSO_SHAPE.ROUNDED_RECTANGLE if rounded else MSO_SHAPE.RECTANGLE
    shape = slide.shapes.add_shape(
        shape_type, Inches(x), Inches(y), Inches(w), Inches(h)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.color.rgb = line if line else fill
    shape.line.width = Pt(0.7 if line else 0.1)
    if rounded:
        shape.adjustments[0] = 0.08
    return shape


def add_circle(slide, x, y, d, fill):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.OVAL, Inches(x), Inches(y), Inches(d), Inches(d)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.color.rgb = fill
    return shape


def add_line(slide, x1, y1, x2, y2, color=LIGHT, width=1, arrow=False):
    shape = slide.shapes.add_connector(
        MSO_CONNECTOR.STRAIGHT,
        Inches(x1),
        Inches(y1),
        Inches(x2),
        Inches(y2),
    )
    shape.line.color.rgb = color
    shape.line.width = Pt(width)
    if arrow:
        shape.line.end_arrowhead = True
    return shape


def add_header(slide, title, subtitle, logo_blob):
    add_rect(slide, 0.32, 0.39, 0.09, 0.28, RED)
    add_rect(slide, 0.32, 0.67, 0.09, 0.20, BLUE)
    add_text(slide, title, 0.56, 0.30, 8.80, 0.50, 23, INK, True)
    add_text(slide, subtitle, 0.58, 0.86, 9.75, 0.32, 10.5, MID)
    slide.shapes.add_picture(
        BytesIO(logo_blob), Inches(10.08), Inches(0.43), width=Inches(2.84)
    )


def add_footer(slide, section, page):
    add_line(slide, 0.58, 6.94, 12.74, 6.94, LIGHT, 0.7)
    add_text(slide, section, 0.58, 7.00, 2.20, 0.18, 7.5, MID)
    add_text(
        slide,
        f"{page:02d}",
        12.40,
        7.00,
        0.33,
        0.18,
        8,
        RGBColor(170, 187, 205),
        False,
        PP_ALIGN.RIGHT,
    )


def add_step(slide, x, y, w, h, number, title, body, accent, fill=WHITE):
    add_rect(slide, x, y, w, h, fill, LIGHT, rounded=True)
    add_rect(slide, x, y, 0.07, h, accent)
    add_circle(slide, x + 0.20, y + 0.18, 0.44, accent)
    add_text(
        slide,
        number,
        x + 0.20,
        y + 0.27,
        0.44,
        0.17,
        8.5,
        WHITE,
        True,
        PP_ALIGN.CENTER,
    )
    add_text(slide, title, x + 0.80, y + 0.14, w - 1.02, 0.31, 12.5, INK, True)
    add_text(
        slide,
        body,
        x + 0.25,
        y + 0.62,
        w - 0.50,
        h - 0.72,
        8.6,
        MID,
        False,
        valign=MSO_ANCHOR.TOP,
    )


def add_data_type(slide, x, y, w, title, examples, accent, number):
    add_rect(slide, x, y, w, 0.90, WHITE, LIGHT, rounded=True)
    add_circle(slide, x + 0.18, y + 0.18, 0.42, accent)
    add_text(
        slide,
        number,
        x + 0.18,
        y + 0.27,
        0.42,
        0.16,
        8,
        WHITE,
        True,
        PP_ALIGN.CENTER,
    )
    add_text(slide, title, x + 0.76, y + 0.13, w - 0.96, 0.28, 10.5, INK, True)
    add_text(
        slide,
        examples,
        x + 0.76,
        y + 0.47,
        w - 0.96,
        0.24,
        7.8,
        MID,
        valign=MSO_ANCHOR.TOP,
    )


def set_notes(slide, text):
    notes_slide = slide.notes_slide
    for shape in notes_slide.shapes:
        if not getattr(shape, "is_placeholder", False):
            continue
        try:
            if shape.placeholder_format.type == 2:
                shape.text_frame.clear()
                shape.text_frame.paragraphs[0].text = text
                return
        except (ValueError, AttributeError):
            continue


def build_model_slide(slide, logo_blob):
    remove_all_shapes(slide)
    set_background(slide)
    add_header(
        slide,
        "城市时空基座大模型研发",
        "基于开源视觉 Qwen-VL，面向地理空间与时间维度约束开展领域适配与能力增强",
        logo_blob,
    )

    add_rect(slide, 0.58, 1.39, 12.15, 0.72, PALE_BLUE, PALE_BLUE, rounded=True)
    add_text(slide, "研发目标", 0.84, 1.55, 0.92, 0.28, 11, BLUE, True)
    add_text(
        slide,
        "构建能够理解城市基础设施对象、空间关系与时间演化信息的多模态时空基座模型，支撑数据理解、空间问答和专题分析。",
        1.88,
        1.50,
        10.45,
        0.36,
        10.3,
        INK,
        True,
    )

    steps = [
        (
            "01",
            "多模态数据构建",
            "组织遥感影像、地图、GIS 矢量、三维场景、业务文本与时序记录，形成面向城市基础设施的训练与评测样本。",
            BLUE,
        ),
        (
            "02",
            "空间约束微调",
            "将坐标、距离、方向、邻接、包含、连通和上下游等空间关系转化为训练约束，增强模型空间关系理解。",
            CYAN,
        ),
        (
            "03",
            "时间约束微调",
            "引入时间点、时间区间、先后顺序、状态变化和历史版本等信息，增强模型对时序事件与对象演化的理解。",
            GREEN,
        ),
        (
            "04",
            "能力评测与应用",
            "围绕对象识别、空间关系问答、时序分析、异常解释和专题辅助开展评测，并形成可调用的模型服务能力。",
            ORANGE,
        ),
    ]
    x_positions = [0.58, 3.67, 6.76, 9.85]
    for step, x in zip(steps, x_positions):
        add_step(slide, x, 2.47, 2.72, 2.45, *step)
    for x in [3.44, 6.53, 9.62]:
        add_text(slide, "→", x, 3.43, 0.18, 0.28, 17, BLUE, True, PP_ALIGN.CENTER)

    add_rect(slide, 0.58, 5.28, 12.15, 0.78, WHITE, LIGHT, rounded=True)
    add_text(slide, "重点研究问题", 0.84, 5.48, 1.32, 0.27, 10.5, NAVY, True)
    research = [
        ("空间关系如何进入模型", "坐标数值、拓扑关系与视觉特征协同表达"),
        ("时间信息如何保持一致", "多源时间口径、状态变化与历史版本统一组织"),
        ("模型结果如何可验证", "规则校验、数据库检索与专家复核共同约束"),
    ]
    for index, (title, body) in enumerate(research):
        x = 2.31 + index * 3.42
        add_text(slide, title, x, 5.40, 1.85, 0.23, 9, BLUE if index == 0 else GREEN if index == 1 else ORANGE, True)
        add_text(slide, body, x, 5.68, 2.90, 0.25, 7.8, MID)

    add_rect(slide, 0.58, 6.27, 12.15, 0.48, NAVY, NAVY, rounded=True)
    add_text(
        slide,
        "形成“视觉理解 + 空间推理 + 时间分析”的城市时空多模态智能能力。",
        0.85,
        6.36,
        11.60,
        0.27,
        10.8,
        WHITE,
        True,
        PP_ALIGN.CENTER,
    )
    add_footer(slide, "后续重点工作（一）", 9)
    set_notes(
        slide,
        (
            "第一项重点工作是城市时空基座大模型研发。技术上以开源视觉Qwen-VL模型为基础，"
            "不是简单进行通用问答，而是针对城市基础设施数据的空间与时间特征开展约束微调。"
            "首先组织遥感影像、地图、GIS矢量、三维场景、业务文本和时序记录等多模态样本；"
            "空间维度重点注入坐标、距离、方向、邻接、包含、连通和上下游等关系；"
            "时间维度重点表达时间点、时间区间、先后顺序、状态变化和历史版本。"
            "最终围绕对象识别、空间关系问答、时序分析、异常解释和专题辅助开展评测，"
            "形成视觉理解、空间推理和时间分析相结合的城市时空智能能力。"
        ),
    )


def build_database_slide(slide, logo_blob):
    remove_all_shapes(slide)
    set_background(slide)
    add_header(
        slide,
        "超融合数据库研发",
        "基于开源 ClickHouse 平台，构建不少于 6 类多模态时空数据的高效存储与检索能力",
        logo_blob,
    )

    add_rect(slide, 0.58, 1.39, 12.15, 0.72, PALE_GREEN, PALE_GREEN, rounded=True)
    add_text(slide, "研发目标", 0.84, 1.55, 0.92, 0.28, 11, GREEN, True)
    add_text(
        slide,
        "围绕统一城市基础设施对象，建立多模态数据接入、时空组织、复合索引与统一检索能力，支撑建模、大模型和数字孪生应用。",
        1.88,
        1.50,
        10.45,
        0.36,
        10.3,
        INK,
        True,
    )

    add_text(slide, "多模态时空数据范围（不少于 6 类）", 0.58, 2.38, 4.10, 0.28, 12.5, INK, True)
    data_types = [
        ("GIS 矢量", "道路、轨道、管线、站点", BLUE),
        ("栅格影像", "遥感、航拍、专题栅格", CYAN),
        ("三维模型", "CIM、BIM、点云与网格", GREEN),
        ("时序数据", "状态、监测、事件与历史", ORANGE),
        ("业务表格", "台账、属性、指标与统计", RED),
        ("文本与文档", "规范、报告、图纸说明", NAVY),
    ]
    positions = [
        (0.58, 2.82),
        (3.18, 2.82),
        (5.78, 2.82),
        (0.58, 3.88),
        (3.18, 3.88),
        (5.78, 3.88),
    ]
    for index, ((title, examples, accent), (x, y)) in enumerate(zip(data_types, positions), start=1):
        add_data_type(slide, x, y, 2.35, title, examples, accent, f"{index:02d}")

    add_rect(slide, 8.46, 2.38, 4.27, 2.40, WHITE, LIGHT, rounded=True)
    add_text(slide, "ClickHouse 核心能力设计", 8.75, 2.60, 3.45, 0.30, 13, INK, True)
    capabilities = [
        ("统一对象主键", "关联对象、空间、时间、来源和版本"),
        ("时空复合组织", "空间分区、时间分区与属性列式存储"),
        ("高效检索分析", "条件过滤、聚合统计和多维组合查询"),
        ("多模态关联", "模型、影像、表格和文档的索引与引用"),
    ]
    for index, (title, body) in enumerate(capabilities):
        y = 3.06 + index * 0.40
        add_circle(slide, 8.76, y + 0.04, 0.16, GREEN)
        add_text(slide, title, 9.05, y - 0.03, 1.25, 0.23, 8.7, GREEN, True)
        add_text(slide, body, 10.32, y - 0.03, 2.02, 0.23, 7.7, MID)

    add_text(slide, "技术路线", 0.58, 5.17, 0.90, 0.27, 11, INK, True)
    route = [
        ("多源接入", "解析与映射"),
        ("对象统一", "编码与分类"),
        ("时空组织", "分区与版本"),
        ("复合索引", "空间/时间/属性"),
        ("统一服务", "检索与分析接口"),
    ]
    for index, (title, body) in enumerate(route):
        x = 1.65 + index * 2.18
        add_rect(slide, x, 5.08, 1.78, 0.76, PALE_BLUE if index < 3 else PALE_GREEN, LIGHT, rounded=True)
        add_text(slide, title, x + 0.12, 5.18, 1.54, 0.22, 9.5, INK, True, PP_ALIGN.CENTER)
        add_text(slide, body, x + 0.12, 5.49, 1.54, 0.18, 7.5, MID, False, PP_ALIGN.CENTER)
        if index < len(route) - 1:
            add_text(slide, "→", x + 1.83, 5.30, 0.27, 0.24, 14, BLUE, True, PP_ALIGN.CENTER)

    add_rect(slide, 0.58, 6.27, 12.15, 0.48, NAVY, NAVY, rounded=True)
    add_text(
        slide,
        "实现多模态时空数据统一组织、高效检索、关联分析与跨专题复用。",
        0.85,
        6.36,
        11.60,
        0.27,
        10.8,
        WHITE,
        True,
        PP_ALIGN.CENTER,
    )
    add_footer(slide, "后续重点工作（二）", 10)
    set_notes(
        slide,
        (
            "第二项重点工作是超融合数据库研发。拟基于开源ClickHouse平台，"
            "面向不少于六类多模态时空数据建立高效存储与检索能力。"
            "数据范围包括GIS矢量、栅格影像、三维模型、时序数据、业务表格以及文本和文档。"
            "研发重点不是简单把不同文件集中存放，而是围绕统一的城市基础设施对象，"
            "建立对象主键、空间位置、时间范围、数据来源和历史版本之间的关联。"
            "在ClickHouse中研究空间分区、时间分区、列式存储和复合查询机制，"
            "形成多源接入、对象统一、时空组织、复合索引和统一服务的完整链路，"
            "为自动化建模、城市时空基座大模型和数字孪生专题提供统一数据底座。"
        ),
    )


def build_ppt():
    presentation = Presentation(SOURCE)
    if len(presentation.slides) < 11:
        raise RuntimeError("当前 PPT 页数不足，无法按既定第9、10页进行替换。")
    logo_blob = presentation.slides[7].shapes[1].image.blob
    build_model_slide(presentation.slides[8], logo_blob)
    build_database_slide(presentation.slides[9], logo_blob)
    presentation.save(OUTPUT)


def font(size, bold=False):
    regular = r"C:\Windows\Fonts\msyh.ttc"
    bold_path = r"C:\Windows\Fonts\msyhbd.ttc"
    return ImageFont.truetype(bold_path if bold else regular, size)


def draw_header(draw, title, subtitle):
    draw.rectangle((38, 48, 48, 82), fill=(235, 35, 58))
    draw.rectangle((38, 82, 48, 106), fill=(16, 112, 224))
    draw.text((68, 42), title, fill=(30, 49, 72), font=font(38, True))
    draw.text((70, 107), subtitle, fill=(91, 111, 135), font=font(19))


def rounded(draw, xy, fill, outline=(215, 228, 240), radius=14, width=2):
    draw.rounded_rectangle(xy, radius=radius, fill=fill, outline=outline, width=width)


def build_model_preview():
    image = Image.new("RGB", (1600, 900), "white")
    draw = ImageDraw.Draw(image)
    draw_header(
        draw,
        "城市时空基座大模型研发",
        "基于开源视觉 Qwen-VL，面向地理空间与时间维度约束开展领域适配与能力增强",
    )
    rounded(draw, (70, 165, 1530, 245), (239, 247, 255), (239, 247, 255))
    draw.text((100, 190), "研发目标", fill=(16, 112, 224), font=font(20, True))
    draw.text(
        (230, 190),
        "构建能够理解城市基础设施对象、空间关系与时间演化信息的多模态时空基座模型",
        fill=(30, 49, 72),
        font=font(19, True),
    )
    steps = [
        ("01", "多模态数据构建", "遥感、地图、GIS、三维场景\n业务文本与时序记录", (16, 112, 224)),
        ("02", "空间约束微调", "坐标、距离、方向、邻接\n包含、连通和上下游", (24, 166, 207)),
        ("03", "时间约束微调", "时间点、时间区间、先后顺序\n状态变化和历史版本", (33, 145, 103)),
        ("04", "能力评测与应用", "对象识别、空间问答、时序分析\n异常解释和专题辅助", (229, 139, 48)),
    ]
    for index, (num, title, body, accent) in enumerate(steps):
        x = 70 + index * 375
        rounded(draw, (x, 290, x + 330, 555), "white")
        draw.rectangle((x, 300, x + 9, 545), fill=accent)
        draw.ellipse((x + 25, 315, x + 80, 370), fill=accent)
        draw.text((x + 52, 342), num, fill="white", font=font(16, True), anchor="mm")
        draw.text((x + 102, 320), title, fill=(30, 49, 72), font=font(22, True))
        draw.multiline_text((x + 30, 400), body, fill=(91, 111, 135), font=font(17), spacing=10)
        if index < 3:
            draw.text((x + 345, 410), "→", fill=(16, 112, 224), font=font(28, True))
    rounded(draw, (70, 610, 1530, 700), "white")
    draw.text((100, 635), "重点研究问题", fill=(16, 50, 88), font=font(20, True))
    draw.text((310, 628), "空间关系如何进入模型", fill=(16, 112, 224), font=font(18, True))
    draw.text((660, 628), "时间信息如何保持一致", fill=(33, 145, 103), font=font(18, True))
    draw.text((1030, 628), "模型结果如何可验证", fill=(229, 139, 48), font=font(18, True))
    rounded(draw, (70, 750, 1530, 815), (16, 50, 88), (16, 50, 88))
    draw.text(
        (800, 782),
        "形成“视觉理解 + 空间推理 + 时间分析”的城市时空多模态智能能力",
        fill="white",
        font=font(21, True),
        anchor="mm",
    )
    draw.line((70, 855, 1530, 855), fill=(215, 228, 240), width=2)
    draw.text((70, 866), "后续重点工作（一）", fill=(91, 111, 135), font=font(13))
    draw.text((1505, 866), "09", fill=(170, 187, 205), font=font(14), anchor="ra")
    PREVIEW_9.parent.mkdir(parents=True, exist_ok=True)
    image.save(PREVIEW_9)


def build_database_preview():
    image = Image.new("RGB", (1600, 900), "white")
    draw = ImageDraw.Draw(image)
    draw_header(
        draw,
        "超融合数据库研发",
        "基于开源 ClickHouse 平台，构建不少于 6 类多模态时空数据的高效存储与检索能力",
    )
    rounded(draw, (70, 165, 1530, 245), (236, 248, 243), (236, 248, 243))
    draw.text((100, 190), "研发目标", fill=(33, 145, 103), font=font(20, True))
    draw.text(
        (230, 190),
        "围绕统一基础设施对象，建立多模态数据接入、时空组织、复合索引与统一检索能力",
        fill=(30, 49, 72),
        font=font(19, True),
    )
    draw.text((70, 280), "多模态时空数据范围（不少于 6 类）", fill=(30, 49, 72), font=font(22, True))
    data = [
        ("GIS 矢量", "道路、轨道、管线、站点", (16, 112, 224)),
        ("栅格影像", "遥感、航拍、专题栅格", (24, 166, 207)),
        ("三维模型", "CIM、BIM、点云与网格", (33, 145, 103)),
        ("时序数据", "状态、监测、事件与历史", (229, 139, 48)),
        ("业务表格", "台账、属性、指标与统计", (235, 35, 58)),
        ("文本与文档", "规范、报告、图纸说明", (16, 50, 88)),
    ]
    for index, (title, body, accent) in enumerate(data):
        col, row = index % 3, index // 3
        x, y = 70 + col * 300, 330 + row * 115
        rounded(draw, (x, y, x + 270, y + 95), "white")
        draw.ellipse((x + 18, y + 18, x + 62, y + 62), fill=accent)
        draw.text((x + 40, y + 40), f"{index + 1:02d}", fill="white", font=font(13, True), anchor="mm")
        draw.text((x + 82, y + 15), title, fill=(30, 49, 72), font=font(18, True))
        draw.text((x + 82, y + 50), body, fill=(91, 111, 135), font=font(13))
    rounded(draw, (1000, 280, 1530, 550), "white")
    draw.text((1035, 310), "ClickHouse 核心能力设计", fill=(30, 49, 72), font=font(22, True))
    caps = [
        "统一对象主键：对象、空间、时间、来源与版本",
        "时空复合组织：空间分区、时间分区、列式存储",
        "高效检索分析：条件过滤、聚合统计与组合查询",
        "多模态关联：模型、影像、表格与文档索引",
    ]
    for index, cap in enumerate(caps):
        draw.ellipse((1038, 365 + index * 42, 1050, 377 + index * 42), fill=(33, 145, 103))
        draw.text((1065, 355 + index * 42), cap, fill=(91, 111, 135), font=font(15))
    draw.text((70, 610), "技术路线", fill=(30, 49, 72), font=font(20, True))
    route = ["多源接入", "对象统一", "时空组织", "复合索引", "统一服务"]
    for index, title in enumerate(route):
        x = 190 + index * 270
        rounded(draw, (x, 595, x + 220, 675), (239, 247, 255) if index < 3 else (236, 248, 243))
        draw.text((x + 110, 635), title, fill=(30, 49, 72), font=font(18, True), anchor="mm")
        if index < 4:
            draw.text((x + 235, 635), "→", fill=(16, 112, 224), font=font(24, True), anchor="mm")
    rounded(draw, (70, 750, 1530, 815), (16, 50, 88), (16, 50, 88))
    draw.text(
        (800, 782),
        "实现多模态时空数据统一组织、高效检索、关联分析与跨专题复用",
        fill="white",
        font=font(21, True),
        anchor="mm",
    )
    draw.line((70, 855, 1530, 855), fill=(215, 228, 240), width=2)
    draw.text((70, 866), "后续重点工作（二）", fill=(91, 111, 135), font=font(13))
    draw.text((1505, 866), "10", fill=(170, 187, 205), font=font(14), anchor="ra")
    image.save(PREVIEW_10)


if __name__ == "__main__":
    build_ppt()
    build_model_preview()
    build_database_preview()
    print(OUTPUT)
    print(PREVIEW_9)
    print(PREVIEW_10)
