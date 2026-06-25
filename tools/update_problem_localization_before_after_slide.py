from io import BytesIO
from pathlib import Path
import zipfile

from PIL import Image, ImageDraw, ImageFont
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_CONNECTOR, MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt


ROOT = Path(r"C:\Users\22838\Desktop\chk\CIMAgent\cim_road_poc")
SOURCE = Path(r"C:\Users\22838\Desktop\chk\CIMAgent\PPT汇报\chkpptx.pptx")
BEFORE_IMAGE = Path(
    r"C:\Users\22838\AppData\Local\Temp\codex-clipboard-93e914ba-7187-438a-ac00-e9317b774306.png"
)
AFTER_IMAGE = Path(
    r"C:\Users\22838\AppData\Local\Temp\codex-clipboard-eaada0e2-3bc5-43d6-a9e8-46c39863bd78.png"
)
OUTPUT = ROOT / "chkpptx_问题定位前后对照版.pptx"
PREVIEW = ROOT / "output" / "chkpptx_问题定位前后对照_第9页预览.png"

FONT = "Microsoft YaHei"
NAVY = RGBColor(16, 50, 88)
BLUE = RGBColor(16, 112, 224)
CYAN = RGBColor(31, 169, 206)
GREEN = RGBColor(39, 146, 104)
ORANGE = RGBColor(242, 143, 35)
RED = RGBColor(239, 31, 55)
INK = RGBColor(26, 46, 70)
MID = RGBColor(83, 105, 130)
LIGHT = RGBColor(211, 226, 239)
PALE_BLUE = RGBColor(241, 247, 253)
PALE_GREEN = RGBColor(236, 248, 243)
PALE_RED = RGBColor(255, 241, 243)
WHITE = RGBColor(255, 255, 255)


def remove_all_shapes(slide):
    for shape in list(slide.shapes):
        element = shape._element
        element.getparent().remove(element)


def add_text(
    slide,
    text,
    x,
    y,
    w,
    h,
    size=12,
    color=INK,
    bold=False,
    align=PP_ALIGN.LEFT,
    valign=MSO_ANCHOR.MIDDLE,
    margin=0.03,
):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    frame = box.text_frame
    frame.clear()
    frame.word_wrap = True
    frame.margin_left = frame.margin_right = Inches(margin)
    frame.margin_top = frame.margin_bottom = Inches(margin)
    frame.vertical_anchor = valign
    paragraph = frame.paragraphs[0]
    paragraph.text = text
    paragraph.alignment = align
    paragraph.space_after = Pt(0)
    paragraph.line_spacing = 1.02
    run = paragraph.runs[0]
    run.font.name = FONT
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    return box


def add_rect(slide, x, y, w, h, fill, line=None, rounded=False, line_width=0.8):
    shape_type = MSO_SHAPE.ROUNDED_RECTANGLE if rounded else MSO_SHAPE.RECTANGLE
    shape = slide.shapes.add_shape(
        shape_type, Inches(x), Inches(y), Inches(w), Inches(h)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.color.rgb = line if line else fill
    shape.line.width = Pt(line_width if line else 0.1)
    if rounded:
        shape.adjustments[0] = 0.07
    return shape


def add_outline(slide, shape_type, x, y, w, h, color, width=2.4):
    shape = slide.shapes.add_shape(
        shape_type, Inches(x), Inches(y), Inches(w), Inches(h)
    )
    shape.fill.background()
    shape.line.color.rgb = color
    shape.line.width = Pt(width)
    return shape


def add_line(slide, x1, y1, x2, y2, color=LIGHT, width=1):
    connector = slide.shapes.add_connector(
        MSO_CONNECTOR.STRAIGHT,
        Inches(x1),
        Inches(y1),
        Inches(x2),
        Inches(y2),
    )
    connector.line.color.rgb = color
    connector.line.width = Pt(width)
    return connector


def add_picture_contain(slide, image_path, x, y, w, h):
    with Image.open(image_path) as image:
        image_w, image_h = image.size
    scale = min(w / image_w, h / image_h)
    draw_w = image_w * scale
    draw_h = image_h * scale
    draw_x = x + (w - draw_w) / 2
    draw_y = y + (h - draw_h) / 2
    slide.shapes.add_picture(
        str(image_path),
        Inches(draw_x),
        Inches(draw_y),
        width=Inches(draw_w),
        height=Inches(draw_h),
    )
    return draw_x, draw_y, draw_w, draw_h


def add_step_row(slide, y, number, title, body, accent):
    add_rect(slide, 4.46, y, 4.40, 0.61, WHITE, LIGHT, rounded=True, line_width=0.7)
    add_rect(slide, 4.46, y, 0.08, 0.61, accent)
    circle = slide.shapes.add_shape(
        MSO_SHAPE.OVAL, Inches(4.64), Inches(y + 0.10), Inches(0.41), Inches(0.41)
    )
    circle.fill.solid()
    circle.fill.fore_color.rgb = accent
    circle.line.color.rgb = accent
    add_text(
        slide,
        number,
        4.64,
        y + 0.205,
        0.41,
        0.15,
        8.8,
        WHITE,
        True,
        PP_ALIGN.CENTER,
    )
    add_text(slide, title, 5.16, y + 0.09, 0.72, 0.20, 10.3, accent, True)
    add_text(slide, body, 5.16, y + 0.29, 3.48, 0.22, 8.25, INK)


def set_notes(slide, text):
    notes_slide = slide.notes_slide
    for shape in notes_slide.shapes:
        if not getattr(shape, "is_placeholder", False):
            continue
        try:
            if shape.placeholder_format.type == 2:
                shape.text_frame.clear()
                shape.text_frame.paragraphs[0].text = text
                return
        except (ValueError, AttributeError):
            continue
    raise RuntimeError("未找到讲者备注正文占位符")


def build_ppt():
    for path in (SOURCE, BEFORE_IMAGE, AFTER_IMAGE):
        if not path.exists():
            raise FileNotFoundError(path)

    presentation = Presentation(SOURCE)
    if len(presentation.slides) < 9:
        raise RuntimeError("当前 PPT 不包含第 9 页")
    slide = presentation.slides[8]
    logo_blob = presentation.slides[7].shapes[1].image.blob

    remove_all_shapes(slide)
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = WHITE

    # Header
    add_rect(slide, 0.32, 0.39, 0.09, 0.28, RED)
    add_rect(slide, 0.32, 0.67, 0.09, 0.20, BLUE)
    add_text(slide, "反思：道路路口问题的定位与修正", 0.56, 0.30, 7.45, 0.50, 23, INK, True)
    add_text(
        slide,
        "案例：斜交支路接入主路时，路侧构件未按路口边界收口",
        0.58,
        0.86,
        8.80,
        0.30,
        10.5,
        MID,
    )
    slide.shapes.add_picture(
        BytesIO(logo_blob), Inches(10.08), Inches(0.43), width=Inches(2.84)
    )

    # Before panel
    add_rect(slide, 0.56, 1.33, 3.66, 3.83, WHITE, RED, rounded=True, line_width=1.2)
    add_rect(slide, 0.56, 1.33, 3.66, 0.43, PALE_RED, PALE_RED, rounded=True)
    add_text(slide, "问题发现", 0.76, 1.43, 0.82, 0.18, 11.0, RED, True)
    add_text(
        slide,
        "直线延伸造成穿插与硬折",
        1.59,
        1.43,
        2.38,
        0.18,
        9.0,
        INK,
        True,
        PP_ALIGN.RIGHT,
    )
    before_x, before_y, before_w, before_h = add_picture_contain(
        slide, BEFORE_IMAGE, 0.68, 1.83, 3.42, 3.15
    )
    # Mark the conflicting termination/overlap in the real screenshot.
    add_outline(
        slide,
        MSO_SHAPE.OVAL,
        before_x + before_w * 0.22,
        before_y + before_h * 0.31,
        before_w * 0.49,
        before_h * 0.45,
        RED,
        2.5,
    )
    add_rect(slide, 0.78, 4.69, 2.84, 0.25, RED, RED, rounded=True)
    add_text(
        slide,
        "路缘/人行道仍沿原道路方向直线扫掠",
        0.82,
        4.71,
        2.76,
        0.16,
        7.6,
        WHITE,
        True,
        PP_ALIGN.CENTER,
    )

    # Five-step localization process
    add_text(slide, "真实定位过程", 4.52, 1.35, 1.30, 0.26, 11.2, NAVY, True)
    add_text(
        slide,
        "从可见缺陷追溯到几何规则",
        6.26,
        1.38,
        2.48,
        0.20,
        8.3,
        MID,
        False,
        PP_ALIGN.RIGHT,
    )
    steps = [
        ("1", "发现", "斜交支路两侧路侧带横穿主路，出现硬折、重叠与断口", BLUE),
        ("2", "隔离", "单独导出问题路口 OBJ/FBX，逐层隐藏后锁定路侧构件", CYAN),
        ("3", "解释", "道路面已连通；路侧带仍独立直扫，端部继续采用直线裁剪", ORANGE),
        ("4", "修正", "依据 arms 方向与断面宽度，生成变宽曲线连接并统一裁剪", RED),
        ("5", "检查", "单路口前后对照，再回归全量模型及 CIM3/CIM4 共用逻辑", GREEN),
    ]
    for index, step in enumerate(steps):
        add_step_row(slide, 1.70 + index * 0.68, *step)

    # After panel
    add_rect(slide, 9.10, 1.33, 3.66, 3.83, WHITE, GREEN, rounded=True, line_width=1.2)
    add_rect(slide, 9.10, 1.33, 3.66, 0.43, PALE_GREEN, PALE_GREEN, rounded=True)
    add_text(slide, "修正结果", 9.30, 1.43, 0.82, 0.18, 11.0, GREEN, True)
    add_text(
        slide,
        "按路口边界生成曲线连接",
        10.13,
        1.43,
        2.38,
        0.18,
        9.0,
        INK,
        True,
        PP_ALIGN.RIGHT,
    )
    after_x, after_y, after_w, after_h = add_picture_contain(
        slide, AFTER_IMAGE, 9.22, 1.83, 3.42, 3.15
    )
    add_outline(
        slide,
        MSO_SHAPE.OVAL,
        after_x + after_w * 0.20,
        after_y + after_h * 0.20,
        after_w * 0.60,
        after_h * 0.58,
        GREEN,
        2.5,
    )
    add_rect(slide, 9.50, 4.69, 2.84, 0.25, GREEN, GREEN, rounded=True)
    add_text(
        slide,
        "曲线连续，路侧构件不再侵入主路",
        9.54,
        4.71,
        2.76,
        0.16,
        7.6,
        WHITE,
        True,
        PP_ALIGN.CENTER,
    )

    # Bottom synthesis
    add_rect(slide, 0.58, 5.44, 5.90, 0.82, PALE_RED, PALE_RED, rounded=True)
    add_text(slide, "问题本质", 0.82, 5.59, 0.80, 0.24, 10.5, RED, True)
    add_text(
        slide,
        "不是参数偏差，而是“直线道路断面规则”没有覆盖“斜交路口连接规则”。",
        1.73,
        5.54,
        4.45,
        0.33,
        9.3,
        INK,
        True,
    )
    add_rect(slide, 6.66, 5.44, 6.07, 0.82, PALE_GREEN, PALE_GREEN, rounded=True)
    add_text(slide, "形成能力", 6.90, 5.59, 0.80, 0.24, 10.5, GREEN, True)
    add_text(
        slide,
        "将夹角、断面宽度和构件类型纳入统一约束，沉淀为可复用的路口连接规则。",
        7.81,
        5.54,
        4.62,
        0.33,
        9.3,
        INK,
        True,
    )
    add_rect(slide, 0.58, 6.43, 12.15, 0.47, NAVY, NAVY, rounded=True)
    add_text(
        slide,
        "定位方法：对象层隔离 → 几何层比对 → 规则层解释 → 单路口验证 → 全量模型回归",
        0.82,
        6.52,
        11.66,
        0.21,
        10.1,
        WHITE,
        True,
        PP_ALIGN.CENTER,
    )

    add_line(slide, 0.58, 7.04, 12.74, 7.04, LIGHT, 0.7)
    add_text(slide, "成果与反思", 0.58, 7.09, 1.50, 0.18, 7.5, MID)
    add_text(
        slide,
        "10",
        12.40,
        7.09,
        0.33,
        0.18,
        8,
        RGBColor(170, 187, 205),
        False,
        PP_ALIGN.RIGHT,
    )

    notes = (
        "这一页不讲抽象的调试流程，而是用这次真实的斜交路口问题说明我们如何完成定位。\n\n"
        "第一张图是问题出现时的模型。可以看到，主路的道路面和支路的道路面已经能够连通，"
        "但支路两侧的人行道、路缘等路侧构件仍按直线道路的方式向前延伸。"
        "在斜交接入条件下，这些构件直接横穿主路，形成明显的硬折、重叠和边界断口。"
        "因此，最初不能把问题笼统归结为道路没有生成成功，而要判断到底是道路面、标线、设施，"
        "还是路侧构件中的哪一层出了问题。\n\n"
        "定位的第一步，是把问题从全量城市模型中隔离出来。我们按路口单独导出 OBJ 和 FBX，"
        "在 Blender 中逐层隐藏标线、路灯等对象。这样可以确认：道路面本身已经连通，"
        "异常集中在 Sidewalk、路缘和相邻路侧带的连接位置。\n\n"
        "继续回查生成逻辑后发现，直线路段上的各类路侧构件，是根据道路中心线和断面宽度独立扫掠生成的。"
        "旧逻辑在进入路口时仍然采用直线端部和固定裁剪范围，没有同时考虑两条路口 arm 的方向夹角、"
        "道路宽度以及构件处于道路外侧的位置关系。这个逻辑在近似正交路口中不容易暴露，"
        "但在斜交支路接入主路时会把原本合理的直线断面延伸成图中的穿插。"
        "所以根因不是简单的坐标偏移或某个宽度参数不准确，而是直线道路规则没有覆盖斜交路口的连接情形。\n\n"
        "修正时，我们没有用局部补片把缺口盖住，而是从路口 arms 和道路断面规则中提取连接点。"
        "根据两条道路的切向方向和各自断面宽度，生成具有连续切向的曲线连接边界，"
        "再构造可变宽的连接带，并通过平面 union、intersection 和 difference 与路口范围统一裁剪。"
        "这样，人行道、路缘和其他路侧构件进入路口时仍保持与直线路段一致的断面语义，"
        "但几何边界会根据斜交角度自然收口。\n\n"
        "第二张图是修正后的结果。两侧边界已经从直线横穿变为连续曲线，支路能够平顺接入主路，"
        "同时没有破坏斑马线和其他附属设施的位置。验证时先比较单路口前后结果，"
        "再回归全量模型，并检查 CIM3 与 CIM4 共用逻辑。"
        "这次工作的价值，不只是修复一个路口，而是把夹角、断面宽度和构件类型沉淀成可复用的路口连接规则。\n\n"
        "汇报时建议先指第一张图中的红圈，说清楚‘道路面正常、路侧构件异常’；"
        "再沿中间五步讲清隔离和根因；最后指第二张图，强调修复结果和规则沉淀。"
    )
    set_notes(slide, notes)

    presentation.save(OUTPUT)
    with zipfile.ZipFile(OUTPUT) as archive:
        if archive.testzip() is not None:
            raise RuntimeError("生成的 PPTX 压缩包校验失败")
    return presentation


def font(size, bold=False):
    regular = r"C:\Windows\Fonts\msyh.ttc"
    bold_path = r"C:\Windows\Fonts\msyhbd.ttc"
    return ImageFont.truetype(bold_path if bold else regular, size)


def rounded(draw, xy, fill, outline=None, radius=16, width=2):
    draw.rounded_rectangle(
        xy,
        radius=radius,
        fill=fill,
        outline=outline if outline else fill,
        width=width,
    )


def contain_pil(image, box):
    x0, y0, x1, y1 = box
    target_w = x1 - x0
    target_h = y1 - y0
    copy = image.convert("RGB")
    copy.thumbnail((target_w, target_h), Image.Resampling.LANCZOS)
    x = x0 + (target_w - copy.width) // 2
    y = y0 + (target_h - copy.height) // 2
    return copy, (x, y)


def build_preview():
    canvas = Image.new("RGB", (1600, 900), "white")
    draw = ImageDraw.Draw(canvas)
    ink = (26, 46, 70)
    mid = (83, 105, 130)
    red = (239, 31, 55)
    blue = (16, 112, 224)
    green = (39, 146, 104)
    navy = (16, 50, 88)
    light = (211, 226, 239)

    draw.rectangle((38, 48, 49, 82), fill=red)
    draw.rectangle((38, 82, 49, 107), fill=blue)
    draw.text((68, 42), "反思：道路路口问题的定位与修正", fill=ink, font=font(38, True))
    draw.text(
        (70, 107),
        "案例：斜交支路接入主路时，路侧构件未按路口边界收口",
        fill=mid,
        font=font(18),
    )

    # Before / after cards
    rounded(draw, (68, 160, 506, 617), "white", red, 16, 3)
    rounded(draw, (68, 160, 506, 211), (255, 241, 243), (255, 241, 243), 16, 1)
    draw.text((92, 174), "问题发现", fill=red, font=font(19, True))
    draw.text((484, 176), "直线延伸造成穿插与硬折", fill=ink, font=font(14, True), anchor="ra")

    rounded(draw, (1093, 160, 1531, 617), "white", green, 16, 3)
    rounded(draw, (1093, 160, 1531, 211), (236, 248, 243), (236, 248, 243), 16, 1)
    draw.text((1117, 174), "修正结果", fill=green, font=font(19, True))
    draw.text((1509, 176), "按路口边界生成曲线连接", fill=ink, font=font(14, True), anchor="ra")

    with Image.open(BEFORE_IMAGE) as source_before:
        before, before_pos = contain_pil(source_before, (83, 222, 491, 593))
    canvas.paste(before, before_pos)
    bx, by = before_pos
    draw.ellipse(
        (
            bx + int(before.width * 0.22),
            by + int(before.height * 0.31),
            bx + int(before.width * 0.71),
            by + int(before.height * 0.76),
        ),
        outline=red,
        width=5,
    )
    rounded(draw, (118, 565, 456, 596), red, red, 11, 1)
    draw.text(
        (287, 580),
        "路侧构件仍按直线扫掠",
        fill="white",
        font=font(13, True),
        anchor="mm",
    )

    with Image.open(AFTER_IMAGE) as source_after:
        after, after_pos = contain_pil(source_after, (1108, 222, 1516, 593))
    canvas.paste(after, after_pos)
    ax, ay = after_pos
    draw.ellipse(
        (
            ax + int(after.width * 0.20),
            ay + int(after.height * 0.20),
            ax + int(after.width * 0.80),
            ay + int(after.height * 0.78),
        ),
        outline=green,
        width=5,
    )
    rounded(draw, (1143, 565, 1481, 596), green, green, 11, 1)
    draw.text(
        (1312, 580),
        "曲线连续，不再侵入主路",
        fill="white",
        font=font(13, True),
        anchor="mm",
    )

    draw.text((539, 163), "真实定位过程", fill=navy, font=font(20, True))
    draw.text((1060, 167), "从可见缺陷追溯到几何规则", fill=mid, font=font(13), anchor="ra")
    steps = [
        ("1", "发现", "斜交支路路侧带横穿主路", blue),
        ("2", "隔离", "单路口导出，逐层锁定路侧构件", (31, 169, 206)),
        ("3", "解释", "道路面正常；路侧带仍独立直扫", (242, 143, 35)),
        ("4", "修正", "按方向与断面宽度生成曲线连接", red),
        ("5", "检查", "单路口对照后回归全量模型", green),
    ]
    for i, (num, title, body, accent) in enumerate(steps):
        top = 205 + i * 76
        rounded(draw, (534, top, 1064, top + 66), "white", light, 12, 2)
        draw.rectangle((534, top + 5, 542, top + 61), fill=accent)
        draw.ellipse((555, top + 13, 598, top + 56), fill=accent)
        draw.text((576, top + 34), num, fill="white", font=font(14, True), anchor="mm")
        draw.text((618, top + 10), title, fill=accent, font=font(16, True))
        draw.text((618, top + 35), body, fill=ink, font=font(13))

    rounded(draw, (70, 650, 778, 745), (255, 241, 243), (255, 241, 243), 15, 1)
    draw.text((96, 670), "问题本质", fill=red, font=font(18, True))
    draw.text((206, 671), "直线道路断面规则没有覆盖斜交路口连接规则", fill=ink, font=font(16, True))
    rounded(draw, (801, 650, 1530, 745), (236, 248, 243), (236, 248, 243), 15, 1)
    draw.text((827, 670), "形成能力", fill=green, font=font(18, True))
    draw.text((937, 671), "将角度、宽度和构件类型沉淀为统一约束", fill=ink, font=font(16, True))
    rounded(draw, (70, 770, 1530, 827), navy, navy, 14, 1)
    draw.text(
        (800, 798),
        "对象层隔离 → 几何层比对 → 规则层解释 → 单路口验证 → 全量模型回归",
        fill="white",
        font=font(19, True),
        anchor="mm",
    )
    draw.line((70, 856, 1530, 856), fill=light, width=2)
    draw.text((70, 866), "成果与反思", fill=mid, font=font(13))
    draw.text((1510, 866), "10", fill=(170, 187, 205), font=font(14), anchor="ra")
    PREVIEW.parent.mkdir(parents=True, exist_ok=True)
    canvas.save(PREVIEW)


if __name__ == "__main__":
    presentation = build_ppt()
    build_preview()
    check = Presentation(OUTPUT)
    notes_text = check.slides[8].notes_slide.notes_text_frame.text
    print(f"output={OUTPUT}")
    print(f"preview={PREVIEW}")
    print(f"slides={len(check.slides)}")
    print(f"slide9_shapes={len(check.slides[8].shapes)}")
    print(f"slide9_notes_chars={len(notes_text)}")
