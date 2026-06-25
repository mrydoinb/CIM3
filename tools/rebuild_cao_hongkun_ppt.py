from pathlib import Path

from PIL import Image
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_CONNECTOR, MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt


ROOT = Path(r"C:\Users\22838\Desktop\chk\CIMAgent\cim_road_poc")
MEDIA = ROOT / "output" / "ppt_unpacked" / "ppt" / "media"
OUT = ROOT / "曹宏坤_技术路线修改版_未来规划扩展版.pptx"

BLUE = RGBColor(0, 112, 224)
DEEP_BLUE = RGBColor(0, 64, 128)
CYAN = RGBColor(36, 180, 225)
PALE_BLUE = RGBColor(238, 247, 255)
INK = RGBColor(31, 45, 61)
MID = RGBColor(91, 107, 126)
LIGHT = RGBColor(220, 230, 240)
WHITE = RGBColor(255, 255, 255)
RED = RGBColor(235, 35, 58)
GREEN = RGBColor(46, 142, 96)
ORANGE = RGBColor(239, 144, 44)
FONT = "Microsoft YaHei"


def add_text(
    slide,
    text,
    x,
    y,
    w,
    h,
    size=18,
    color=INK,
    bold=False,
    align=PP_ALIGN.LEFT,
    valign=MSO_ANCHOR.TOP,
    margin=0.04,
    linesp=1.0,
):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.clear()
    tf.margin_left = tf.margin_right = Inches(margin)
    tf.margin_top = tf.margin_bottom = Inches(margin)
    tf.vertical_anchor = valign
    p = tf.paragraphs[0]
    p.text = text
    p.alignment = align
    p.line_spacing = linesp
    r = p.runs[0]
    r.font.name = FONT
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.color.rgb = color
    return box


def add_rect(slide, x, y, w, h, fill=WHITE, line=None, radius=False):
    shape_type = MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE
    shp = slide.shapes.add_shape(shape_type, Inches(x), Inches(y), Inches(w), Inches(h))
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill
    shp.line.color.rgb = line if line else fill
    if radius:
        shp.adjustments[0] = 0.08
    return shp


def add_line(slide, x1, y1, x2, y2, color=BLUE, width=1.5):
    line = slide.shapes.add_connector(
        MSO_CONNECTOR.STRAIGHT, Inches(x1), Inches(y1), Inches(x2), Inches(y2)
    )
    line.line.color.rgb = color
    line.line.width = Pt(width)
    return line


def add_circle(slide, x, y, d, fill=BLUE):
    shp = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(x), Inches(y), Inches(d), Inches(d))
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill
    shp.line.color.rgb = fill
    return shp


def add_chevron(slide, x, y, w, h, fill=BLUE):
    shp = slide.shapes.add_shape(MSO_SHAPE.CHEVRON, Inches(x), Inches(y), Inches(w), Inches(h))
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill
    shp.line.color.rgb = fill
    return shp


def add_picture_crop(slide, path, x, y, w, h):
    with Image.open(path) as im:
        iw, ih = im.size
    image_ratio = iw / ih
    box_ratio = w / h
    pic = slide.shapes.add_picture(str(path), Inches(x), Inches(y), width=Inches(w), height=Inches(h))
    if image_ratio > box_ratio:
        visible = box_ratio / image_ratio
        pic.crop_left = pic.crop_right = (1 - visible) / 2
    else:
        visible = image_ratio / box_ratio
        pic.crop_top = pic.crop_bottom = (1 - visible) / 2
    return pic


def add_header(slide, title, section, page, dark=False):
    title_color = WHITE if dark else INK
    sub_color = RGBColor(190, 224, 255) if dark else BLUE
    add_text(slide, section.upper(), 0.55, 0.25, 2.5, 0.25, 8, sub_color, True)
    add_text(slide, title, 0.55, 0.52, 11.65, 0.5, 22, title_color, True)
    add_rect(slide, 0.55, 1.07, 0.46, 0.035, RED)
    add_rect(slide, 1.02, 1.07, 1.2, 0.035, BLUE)
    logo = MEDIA / "image2.png"
    if logo.exists():
        slide.shapes.add_picture(str(logo), Inches(10.72), Inches(0.28), width=Inches(2.05))
    add_text(slide, f"{page:02d}", 12.42, 7.05, 0.35, 0.2, 8, RGBColor(180, 195, 210), align=PP_ALIGN.RIGHT)


def add_footer(slide, text="技术路线 · 问题意识 · 能力沉淀"):
    add_line(slide, 0.55, 6.98, 12.78, 6.98, LIGHT, 0.6)
    add_text(slide, text, 0.55, 7.06, 5.6, 0.18, 7, MID)


def add_card(slide, x, y, w, h, title, body, accent=BLUE, number=None):
    add_rect(slide, x, y, w, h, WHITE, LIGHT, radius=True)
    add_rect(slide, x, y, 0.08, h, accent)
    tx = x + 0.28
    tw = w - 0.55
    if number:
        add_text(slide, number, x + 0.28, y + 0.23, 0.55, 0.35, 16, accent, True)
        tx = x + 0.9
        tw = w - 1.12
    add_text(slide, title, tx, y + 0.2, tw, 0.34, 14.5, INK, True)
    add_text(slide, body, tx, y + 0.66, tw, h - 0.76, 10.6, MID, linesp=1.15)


def new_deck():
    prs = Presentation()
    prs.slide_width = Inches(13.333333)
    prs.slide_height = Inches(7.5)
    return prs


def add_slide(prs, bg=WHITE):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = bg
    return slide


def build():
    prs = new_deck()

    # 1 Cover
    s = add_slide(prs, DEEP_BLUE)
    add_picture_crop(s, MEDIA / "image1.png", 0, 0, 13.333, 7.5)
    overlay = add_rect(s, 0, 0, 13.333, 7.5, DEEP_BLUE)
    overlay.fill.transparency = 26
    add_rect(s, 0.0, 5.92, 8.6, 1.58, RGBColor(0, 74, 150))
    add_rect(s, 0.62, 1.02, 0.12, 3.55, RED)
    add_text(s, "城市基础设施数字化建模与时空数据融合", 1.02, 1.12, 8.8, 0.48, 18, RGBColor(188, 225, 255), True)
    add_text(s, "技术路线与阶段进展", 1.0, 1.78, 8.8, 0.92, 34, WHITE, True)
    add_text(s, "面向超融合数据库与数字孪生专项的共性技术能力建设", 1.02, 3.16, 9.3, 0.5, 16, WHITE)
    add_text(s, "汇报人：曹宏坤", 1.02, 6.22, 3.2, 0.32, 12, WHITE, True)
    add_text(s, "2026.06.23", 1.02, 6.66, 2.2, 0.26, 10, RGBColor(195, 225, 250))
    s.shapes.add_picture(str(MEDIA / "image2.png"), Inches(10.54), Inches(6.73), width=Inches(2.2))

    # 2 Focus
    s = add_slide(prs)
    add_header(s, "汇报重点：围绕总体目标、关键问题与后续重点工作展开", "REPORT FOCUS", 2)
    prompts = [
        ("01", "总体路线", "构建数据治理、规则表达、\n空间建模与应用服务链条。"),
        ("02", "关键问题", "聚焦数据一致性、对象化建模、\n空间关系与数据资产维护。"),
        ("03", "应用支撑", "支撑超融合数据库、数字孪生专项\n及多类基础设施建模。"),
    ]
    for i, (n, t, b) in enumerate(prompts):
        x = 0.72 + i * 4.15
        add_rect(s, x, 1.58, 3.62, 4.54, PALE_BLUE if i != 1 else RGBColor(232, 244, 255), LIGHT, radius=True)
        add_circle(s, x + 0.34, 1.96, 0.72, BLUE if i != 1 else RED)
        add_text(s, n, x + 0.34, 2.13, 0.72, 0.28, 13, WHITE, True, PP_ALIGN.CENTER, MSO_ANCHOR.MIDDLE)
        add_text(s, t, x + 0.35, 2.94, 2.9, 0.42, 20, INK, True)
        add_text(s, b, x + 0.35, 3.64, 2.92, 1.3, 11.2, MID, linesp=1.15)
        add_line(s, x + 0.35, 5.46, x + 2.95, 5.46, BLUE if i != 1 else RED, 2)
    add_footer(s)

    # 3 Capability map
    s = add_slide(prs)
    add_header(s, "阶段定位：以典型场景形成共性技术，支撑城市数字底座建设", "01 / ROUTE", 3)
    items = [
        ("数据接入", "GIS / CAD / 业务表\n多源异构数据治理", BLUE),
        ("规则表达", "专业规则与分级标准\n参数化组织", CYAN),
        ("空间建模", "道路、隧道、管线\n三维对象表达", ORANGE),
        ("语义关联", "对象、属性、关系\n全过程可追溯", GREEN),
        ("应用支撑", "超融合数据库\n数字孪生专项", RED),
    ]
    for i, (t, b, c) in enumerate(items):
        x = 0.55 + i * 2.52
        add_rect(s, x, 2.08, 2.08, 2.18, WHITE, LIGHT, radius=True)
        add_rect(s, x, 2.08, 2.08, 0.4, c)
        add_text(s, t, x + 0.18, 2.76, 1.72, 0.34, 15, INK, True, PP_ALIGN.CENTER)
        add_text(s, b, x + 0.18, 3.28, 1.72, 0.66, 10.8, MID, align=PP_ALIGN.CENTER, linesp=1.15)
        if i < len(items) - 1:
            add_chevron(s, x + 2.15, 2.92, 0.28, 0.48, RGBColor(180, 210, 235))
    add_rect(s, 0.72, 5.38, 11.92, 0.64, PALE_BLUE, LIGHT, radius=True)
    add_text(s, "核心目标：形成多源数据到数字对象的规范化转换能力，实现对象可解释、关系可查询、数据资产可持续演进。",
             0.92, 5.55, 11.5, 0.3, 13, DEEP_BLUE, True, PP_ALIGN.CENTER)
    add_footer(s)

    # 4 Selection
    s = add_slide(prs)
    add_header(s, "技术体系：平台、算法与数据底座协同，形成可持续演进架构", "01 / ROUTE", 4)
    rows = [
        ("三维表达平台", "承担场景组织、三维表达与交互展示", "定位为对象表达与空间关系呈现层，支撑专题应用拓展", BLUE),
        ("空间建模算法", "承担空间构造、分级表达与质量控制", "沉淀可迁移的模型生成能力，服务多类基础设施对象", ORANGE),
        ("时空数据底座", "承担对象、属性、关系、版本与索引管理", "与超融合数据库研发衔接，支撑检索、分析与业务复用", GREEN),
        ("专项应用体系", "承担面向业务场景的成果组织与应用闭环", "支撑数字孪生专项从展示型成果向业务型应用延伸", RED),
    ]
    for i, (name, role, point, color) in enumerate(rows):
        y = 1.48 + i * 1.16
        add_rect(s, 0.72, y, 11.92, 0.92, WHITE if i % 2 == 0 else PALE_BLUE, LIGHT)
        add_rect(s, 0.72, y, 0.08, 0.92, color)
        add_text(s, name, 1.0, y + 0.24, 1.55, 0.3, 13.5, color, True)
        add_text(s, role, 2.82, y + 0.22, 3.3, 0.34, 11.5, INK, True)
        add_text(s, point, 6.38, y + 0.18, 5.72, 0.46, 10.6, MID)
    add_rect(s, 0.72, 6.25, 11.92, 0.5, DEEP_BLUE)
    add_text(s, "总体架构由“数据底座 + 规则算法 + 三维表达 + 专项应用”共同构成，强调能力协同与可持续迭代。",
             0.92, 6.38, 11.5, 0.24, 12.5, WHITE, True, PP_ALIGN.CENTER)
    add_footer(s)

    # 5 Data flow
    s = add_slide(prs)
    add_header(s, "总体路线：以数据链为主线，构建可追溯的城市基础设施对象体系", "01 / ROUTE", 5)
    stages = [
        ("源数据", "GIS / CAD\n业务台账", BLUE),
        ("标准化", "坐标 / 编码\n分类 / 质量", CYAN),
        ("规则层", "分级标准\n专业约束", GREEN),
        ("模型层", "三维对象\n空间关系", ORANGE),
        ("语义层", "对象属性\n关系索引", RED),
        ("服务层", "孪生场景\n查询分析", DEEP_BLUE),
    ]
    for i, (t, b, c) in enumerate(stages):
        x = 0.45 + i * 2.13
        add_rect(s, x, 2.08, 1.78, 2.05, WHITE, LIGHT, radius=True)
        add_rect(s, x, 2.08, 1.78, 0.38, c)
        add_text(s, t, x + 0.12, 2.72, 1.54, 0.35, 14, INK, True, PP_ALIGN.CENTER)
        add_text(s, b, x + 0.14, 3.28, 1.5, 0.56, 10.5, MID, align=PP_ALIGN.CENTER)
        if i < len(stages) - 1:
            add_chevron(s, x + 1.83, 2.88, 0.28, 0.48, RGBColor(180, 210, 235))
    add_card(s, 0.72, 4.82, 5.65, 1.15, "当前阶段", "以 CIM 道路建模为典型场景，打通数据标准化、空间构造和语义关联链路。", BLUE, "A")
    add_card(s, 6.92, 4.82, 5.65, 1.15, "后续阶段", "将阶段形成的方法接入超融合数据库与数字孪生专项，支撑综合型城市基础设施应用。", RED, "B")
    add_footer(s)

    # 6 Current practice
    s = add_slide(prs)
    add_header(s, "阶段实践：以道路建模贯通规则化生成与语义追溯链路", "02 / CURRENT PRACTICE", 6)
    add_picture_crop(s, MEDIA / "image8.png", 7.75, 1.48, 4.88, 3.15)
    add_rect(s, 7.75, 4.33, 4.88, 0.3, DEEP_BLUE)
    add_text(s, "典型道路场景用于贯通数据—模型—语义链路", 7.92, 4.38, 4.5, 0.2, 9.5, WHITE, True, PP_ALIGN.CENTER)
    points = [
        ("输入", "中心线、分类、宽度、断面等基础数据"),
        ("处理", "依据专业规则形成分级模型与对象关系"),
        ("输出", "三维对象、语义数据和典型问题样本库"),
        ("价值", "建立模型对象与数据资产之间的关联机制"),
    ]
    for i, (t, b) in enumerate(points):
        y = 1.55 + i * 1.1
        add_circle(s, 0.78, y, 0.62, BLUE if i < 3 else RED)
        add_text(s, str(i + 1), 0.78, y + 0.14, 0.62, 0.25, 12, WHITE, True, PP_ALIGN.CENTER, MSO_ANCHOR.MIDDLE)
        add_text(s, t, 1.72, y + 0.02, 1.2, 0.32, 15, INK, True)
        add_text(s, b, 2.78, y + 0.02, 3.9, 0.36, 11, MID)
    add_text(s, "该阶段重点在于形成可观察、可复盘、可扩展的技术链路，为多类型设施建模提供依据。",
             0.76, 6.25, 11.8, 0.34, 13, DEEP_BLUE, True, PP_ALIGN.CENTER)
    add_footer(s)

    # 7 Problem 1
    s = add_slide(prs)
    add_header(s, "关键问题一：模型质量受数据一致性、规则表达与空间关系共同影响", "03 / PROBLEMS", 7)
    add_picture_crop(s, MEDIA / "image10.png", 0.72, 1.48, 5.28, 4.55)
    add_rect(s, 0.72, 5.68, 5.28, 0.35, DEEP_BLUE)
    add_text(s, "典型路口场景：用于分析规则冲突与空间连续性", 0.88, 5.75, 4.95, 0.2, 9.5, WHITE, True, PP_ALIGN.CENTER)
    problems = [
        ("数据一致性", "不同来源数据在空间精度、分类字段和属性完备性方面存在差异。", BLUE),
        ("规则可配置性", "专业规则需具备参数化和版本化能力，避免形成一次性项目逻辑。", ORANGE),
        ("对象统一性", "模型展示、语义查询和业务应用需共享统一的对象关系体系。", RED),
    ]
    for i, (t, b, c) in enumerate(problems):
        add_card(s, 6.42, 1.62 + i * 1.34, 6.15, 1.03, t, b, c, f"0{i+1}")
    add_rect(s, 6.42, 5.72, 6.15, 0.56, PALE_BLUE, LIGHT, radius=True)
    add_text(s, "工作重点：由单点模型修正转向数据标准、规则表达和对象关系的系统化治理。",
             6.66, 5.88, 5.7, 0.25, 12, BLUE, True, PP_ALIGN.CENTER)
    add_footer(s)

    # 8 Problem 2
    s = add_slide(prs)
    add_header(s, "关键问题二：专项建设需由文件管理提升为数据资产组织", "03 / PROBLEMS", 8)
    phenomena = [
        ("文件孤立", "仅以文件形态管理，后续追踪、复用和增量更新能力不足。"),
        ("数据分散", "模型、表格、GIS 与业务系统之间缺少统一对象和索引机制。"),
        ("应用分离", "可视化展示、检索分析和业务汇报各自建设，难以形成协同。"),
    ]
    for i, (tag, body) in enumerate(phenomena):
        x = 0.72 + i * 4.03
        add_rect(s, x, 1.72, 3.58, 1.4, WHITE, LIGHT, radius=True)
        add_text(s, tag, x + 0.28, 1.92, 1.4, 0.32, 14, RED if i == 0 else ORANGE, True)
        add_text(s, body, x + 0.28, 2.43, 2.92, 0.48, 9.2, MID, linesp=1.1)
    rules = [
        ("统一编码", "设施对象有稳定 ID"),
        ("统一索引", "空间、时间、属性可检索"),
        ("统一版本", "数据更新可追踪"),
        ("统一接口", "建模、孪生、业务系统可共享"),
    ]
    add_text(s, "数据资产化建设重点", 0.72, 3.64, 2.6, 0.3, 12, BLUE, True)
    for i, (t, b) in enumerate(rules):
        x = 0.72 + i * 3.06
        add_circle(s, x, 4.14, 0.54, BLUE)
        add_text(s, str(i + 1), x, 4.25, 0.54, 0.22, 11, WHITE, True, PP_ALIGN.CENTER, MSO_ANCHOR.MIDDLE)
        add_text(s, t, x + 0.72, 4.12, 1.38, 0.3, 13, INK, True)
        add_text(s, b, x + 0.72, 4.58, 1.82, 0.48, 10.2, MID)
    add_rect(s, 0.72, 5.9, 11.92, 0.54, DEEP_BLUE)
    add_text(s, "超融合数据库研发需重点承接对象化、索引化和版本化能力，使模型对象具备查询、关联和更新基础。", 0.95, 6.05, 11.45, 0.24, 13, WHITE, True, PP_ALIGN.CENTER)
    add_footer(s)

    # 9 Database / twin
    s = add_slide(prs)
    add_header(s, "面向超融合数据库：构建统一的城市基础设施时空数据组织", "04 / NEXT CAPABILITY", 9)
    add_picture_crop(s, MEDIA / "image9.png", 0.72, 1.48, 5.15, 4.58)
    add_rect(s, 0.72, 5.72, 5.15, 0.34, RGBColor(12, 24, 36))
    add_text(s, "城市级对象网络：统一组织空间对象、属性信息与关联关系", 0.9, 5.78, 4.8, 0.2, 9.5, WHITE, True, PP_ALIGN.CENTER)
    add_card(s, 6.32, 1.52, 5.98, 0.95, "对象层", "面向道路、隧道、管线、站点、设施等对象建立统一编码。", BLUE, "01")
    add_card(s, 6.32, 2.7, 5.98, 0.95, "关系层", "表达空间邻接、上下游、穿越、权属和时间版本等关系。", GREEN, "02")
    add_card(s, 6.32, 3.88, 5.98, 0.95, "服务层", "支撑检索、统计、可视化分析和数字孪生应用服务。", RED, "03")
    add_rect(s, 6.32, 5.45, 5.98, 0.58, PALE_BLUE, LIGHT, radius=True)
    add_text(s, "建设重点在于形成对象、关系、索引和持续更新机制，而非仅存储模型文件。",
             6.55, 5.62, 5.55, 0.24, 11.5, DEEP_BLUE, True, PP_ALIGN.CENTER)
    add_footer(s)

    # 10 Digital twin delivery
    s = add_slide(prs)
    add_header(s, "面向数字孪生专项：由场景表达拓展至业务研判与闭环应用", "04 / NEXT CAPABILITY", 10)
    stages = [
        ("现状表达", "基础设施对象\n三维表达", BLUE),
        ("数据融合", "时空数据\n业务数据", CYAN),
        ("关系理解", "空间关系\n运行约束", ORANGE),
        ("场景应用", "查询统计\n研判分析", RED),
        ("反馈更新", "问题回流\n版本迭代", GREEN),
    ]
    for i, (t, b, c) in enumerate(stages):
        x = 0.42 + i * 2.53
        add_circle(s, x + 0.76, 1.62, 0.78, c)
        add_text(s, str(i + 1), x + 0.76, 1.8, 0.78, 0.28, 13, WHITE, True, PP_ALIGN.CENTER, MSO_ANCHOR.MIDDLE)
        add_text(s, t, x + 0.28, 2.72, 1.78, 0.34, 15, INK, True, PP_ALIGN.CENTER)
        add_text(s, b, x + 0.2, 3.26, 1.96, 0.62, 10.5, MID, align=PP_ALIGN.CENTER)
        if i < 4:
            add_chevron(s, x + 2.08, 1.77, 0.32, 0.48, RGBColor(188, 211, 230))
    add_card(s, 0.72, 4.78, 5.65, 1.12, "已形成的基础", "自动建模、对象语义、空间关系和三维表达方法。", BLUE, "A")
    add_card(s, 6.92, 4.78, 5.65, 1.12, "需进一步完善", "数据更新机制、业务指标体系、智能辅助理解与专题应用模板。", RED, "B")
    add_footer(s)

    # 11 Reuse
    s = add_slide(prs)
    add_header(s, "方法复用：形成面向多类城市基础设施对象的共性技术框架", "05 / REFLECTION", 11)
    sequence = ["对象抽象", "数据标准", "规则配置", "模型表达", "应用服务"]
    for i, label in enumerate(sequence):
        x = 0.7 + i * 2.5
        add_rect(s, x, 1.78, 1.82, 0.74, BLUE if i in (0, 4) else WHITE, BLUE, radius=True)
        add_text(s, label, x + 0.12, 2.0, 1.58, 0.28, 12, WHITE if i in (0, 4) else INK, True, PP_ALIGN.CENTER)
        if i < 4:
            add_chevron(s, x + 1.9, 1.92, 0.34, 0.44, RGBColor(180, 210, 235))
    cases = [
        ("道路 / 交通", "研究空间连续、分级表达与对象追溯方法", BLUE),
        ("轨道 / 地下空间", "扩展至区间、站点、构件和空间关系表达", ORANGE),
        ("地下管线", "面向管径、埋深、节点、碰撞与权属关系组织", GREEN),
        ("城市孪生专题", "将基础设施对象接入场景、指标与应用体系", RED),
    ]
    for i, (name, focus, color) in enumerate(cases):
        x = 0.72 + (i % 2) * 6.15
        y = 3.2 + (i // 2) * 1.35
        add_card(s, x, y, 5.68, 1.05, name, focus, color, f"0{i+1}")
    add_text(s, "共性部分在于对象化、规则化和语义化方法；差异部分体现在各专业规则、约束条件和应用场景。",
             0.72, 6.22, 11.92, 0.34, 13, DEEP_BLUE, True, PP_ALIGN.CENTER)
    add_footer(s)

    # 12 Reflection and next steps
    s = add_slide(prs)
    add_header(s, "下一阶段：围绕三条主线推进技术能力体系化建设", "05 / REFLECTION", 12)
    reflections = [
        ("自动化建模", "持续完善道路、轨道、地下空间等对象生成方法，形成可复用建模能力。", BLUE),
        ("超融合数据库", "围绕多模态时空数据，建设统一对象、索引、关系和版本管理能力。", GREEN),
        ("数字孪生专项", "面向专题建设，形成场景模板、指标体系和智能辅助分析能力。", RED),
    ]
    for i, (t, b, c) in enumerate(reflections):
        add_card(s, 0.72, 1.5 + i * 1.35, 6.25, 1.08, t, b, c, f"0{i+1}")
    add_rect(s, 7.34, 1.5, 5.3, 3.84, DEEP_BLUE, DEEP_BLUE, radius=True)
    add_text(s, "后续重点工作", 7.72, 1.78, 4.54, 0.42, 18, WHITE, True)
    nexts = [
        "完善多源数据标准、对象编码与分类体系",
        "建设道路、轨道、管线、地下空间对象模板库",
        "研究规则配置、参数管理与版本演进机制",
        "参与超融合数据库时空索引、关系模型和数据服务设计",
        "面向数字孪生专项建设场景指标、专题图层和智能分析能力",
        "凝练论文、专利、技术资料与可复用示范样例",
    ]
    for i, text in enumerate(nexts):
        y = 2.36 + i * 0.47
        add_circle(s, 7.73, y, 0.28, RED if i == 0 else CYAN)
        add_text(s, str(i + 1), 7.73, y + 0.05, 0.28, 0.16, 7.5, WHITE, True, PP_ALIGN.CENTER, MSO_ANCHOR.MIDDLE)
        add_text(s, text, 8.22, y - 0.02, 3.95, 0.32, 9.1, WHITE)
    add_rect(s, 7.34, 5.64, 5.3, 0.7, PALE_BLUE, LIGHT, radius=True)
    add_text(s, "阶段工作的核心价值在于形成", 7.54, 5.76, 4.92, 0.22, 10.5, DEEP_BLUE, True, PP_ALIGN.CENTER)
    add_text(s, "可迁移、可复用、可持续演进的技术体系。", 7.54, 6.02, 4.92, 0.22, 10.5, DEEP_BLUE, True, PP_ALIGN.CENTER)
    add_footer(s, "从阶段积累走向体系化能力")

    # 13 Future plan: hyper-converged database
    s = add_slide(prs)
    add_header(s, "未来规划一：面向超融合数据库的研发重点", "06 / FUTURE PLAN", 13)
    add_rect(s, 0.72, 1.46, 11.92, 0.78, PALE_BLUE, LIGHT, radius=True)
    add_text(s, "建设目标：形成面向城市基础设施的多模态时空数据底座，支撑模型对象、业务数据与空间关系的统一组织。",
             0.98, 1.66, 11.35, 0.3, 13, DEEP_BLUE, True, PP_ALIGN.CENTER)
    db_items = [
        ("多源接入", "统一接入 GIS、CAD、模型、表格\n与业务台账等多类型数据。", BLUE),
        ("对象建模", "建立设施对象编码、分类体系、\n属性结构和跨专业对象模板。", CYAN),
        ("时空索引", "研究空间、时间、属性\n和语义标签的复合索引。", GREEN),
        ("关系模型", "表达邻接、穿越、上下游、权属、状态和版本等对象关系。", ORANGE),
        ("数据服务", "形成面向建模、孪生、查询、统计和智能分析的数据服务接口。", RED),
    ]
    for i, (title, body, color) in enumerate(db_items):
        x = 0.72 + (i % 3) * 4.05
        y = 2.72 + (i // 3) * 1.55
        w = 3.55 if i < 3 else 5.55
        if i >= 3:
            x = 1.52 + (i - 3) * 5.9
        add_rect(s, x, y, w, 1.18, WHITE, LIGHT, radius=True)
        add_rect(s, x, y, 0.08, 1.18, color)
        add_text(s, f"0{i+1}", x + 0.25, y + 0.22, 0.58, 0.28, 13, color, True)
        add_text(s, title, x + 0.9, y + 0.2, 1.55, 0.3, 13.5, INK, True)
        add_text(s, body, x + 0.9, y + 0.58, w - 1.22, 0.48, 8.9, MID, linesp=1.05)
    add_rect(s, 0.72, 6.18, 11.92, 0.52, DEEP_BLUE)
    add_text(s, "重点不止于数据存储，而是建设“对象可管理、关系可计算、服务可复用”的城市基础设施数据底座。",
             0.95, 6.32, 11.45, 0.22, 12.2, WHITE, True, PP_ALIGN.CENTER)
    add_footer(s, "未来规划 · 超融合数据库研发")

    # 14 Future plan: digital twin special topic
    s = add_slide(prs)
    add_header(s, "未来规划二：数字孪生专项课题交付内容", "06 / FUTURE PLAN", 14)
    add_rect(s, 0.72, 1.42, 5.75, 4.82, PALE_BLUE, LIGHT, radius=True)
    add_text(s, "课题定位", 1.02, 1.72, 1.5, 0.34, 15, BLUE, True)
    add_text(s, "围绕城市基础设施对象，形成可展示、可查询、可分析的专题能力；\n支撑数字孪生场景持续扩展。",
             1.02, 2.14, 4.95, 0.72, 10.8, INK, True, linesp=1.12)
    twin_layers = [
        ("场景底座", "道路、轨道、管线、地下空间等三维对象组织"),
        ("专题图层", "设施分布、空间关系、运行约束和状态信息表达"),
        ("指标体系", "面向管理、研判和汇报的指标口径与展示逻辑"),
        ("智能分析", "结合视觉语言模型与时空数据开展辅助理解和问答"),
    ]
    for i, (title, body) in enumerate(twin_layers):
        y = 3.06 + i * 0.68
        add_circle(s, 1.05, y, 0.34, BLUE if i < 2 else RED)
        add_text(s, str(i + 1), 1.05, y + 0.07, 0.34, 0.15, 8, WHITE, True, PP_ALIGN.CENTER, MSO_ANCHOR.MIDDLE)
        add_text(s, title, 1.52, y - 0.02, 1.1, 0.24, 11.5, INK, True)
        add_text(s, body, 2.58, y - 0.02, 3.45, 0.28, 9.6, MID)
    add_rect(s, 6.88, 1.42, 5.75, 4.82, DEEP_BLUE, DEEP_BLUE, radius=True)
    add_text(s, "交付组织", 7.22, 1.72, 2.0, 0.34, 15, WHITE, True)
    deliverables = [
        ("专题场景", "形成面向汇报和业务分析的孪生场景模板"),
        ("数据关联", "建立模型对象与数据库对象、业务属性的关联机制"),
        ("应用原型", "形成查询、统计、专题展示和智能问答的应用雏形"),
        ("技术沉淀", "同步沉淀接口说明、数据规范、场景配置和演示材料"),
    ]
    for i, (title, body) in enumerate(deliverables):
        y = 2.36 + i * 0.78
        add_circle(s, 7.24, y, 0.32, RED if i == 0 else CYAN)
        add_text(s, str(i + 1), 7.24, y + 0.06, 0.32, 0.15, 8, WHITE, True, PP_ALIGN.CENTER, MSO_ANCHOR.MIDDLE)
        add_text(s, title, 7.72, y - 0.02, 1.2, 0.26, 11.5, WHITE, True)
        add_text(s, body, 8.88, y - 0.02, 3.1, 0.3, 9.3, WHITE)
    add_rect(s, 6.88, 5.45, 5.75, 0.52, RGBColor(0, 86, 160), RGBColor(0, 86, 160), radius=True)
    add_text(s, "目标：支撑专项课题从“场景表达”进一步走向“数据驱动的专题应用”。",
             7.12, 5.6, 5.25, 0.22, 10.5, WHITE, True, PP_ALIGN.CENTER)
    add_footer(s, "未来规划 · 数字孪生专项课题")

    # 15 End
    s = add_slide(prs, DEEP_BLUE)
    add_picture_crop(s, MEDIA / "image11.png", 0, 0, 13.333, 7.5)
    overlay = add_rect(s, 0, 0, 13.333, 7.5, DEEP_BLUE)
    overlay.fill.transparency = 25
    add_text(s, "谢谢", 0.95, 2.12, 3.2, 0.78, 38, WHITE, True)
    add_text(s, "THANK YOU", 0.98, 3.02, 2.4, 0.32, 13, RGBColor(190, 225, 255), True)
    add_rect(s, 0.98, 3.62, 0.52, 0.05, RED)
    add_text(s, "推动建模能力向数据资产能力拓展，", 0.98, 4.04, 6.2, 0.32, 18, WHITE, True)
    add_text(s, "支撑城市数字孪生与专项应用建设。", 0.98, 4.43, 6.2, 0.32, 18, WHITE, True)
    s.shapes.add_picture(str(MEDIA / "image2.png"), Inches(10.48), Inches(6.72), width=Inches(2.28))

    prs.save(OUT)
    print(OUT)


if __name__ == "__main__":
    build()
