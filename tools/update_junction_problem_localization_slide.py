from io import BytesIO
from pathlib import Path

from PIL import Image, ImageDraw, ImageFont
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_CONNECTOR, MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt


ROOT = Path(r"C:\Users\22838\Desktop\chk\CIMAgent\cim_road_poc")
SOURCE = Path(r"C:\Users\22838\Desktop\chk\CIMAgent\PPT汇报\chkpptx.pptx")
OUTPUT = ROOT / "chkpptx_问题定位真实案例版.pptx"
PREVIEW = ROOT / "output" / "chkpptx_问题定位真实案例_第9页预览.png"

FONT = "Microsoft YaHei"

NAVY = RGBColor(16, 50, 88)
BLUE = RGBColor(16, 112, 224)
CYAN = RGBColor(24, 166, 207)
GREEN = RGBColor(33, 145, 103)
ORANGE = RGBColor(242, 143, 35)
RED = RGBColor(239, 31, 55)
INK = RGBColor(26, 46, 70)
MID = RGBColor(86, 108, 134)
LIGHT = RGBColor(212, 227, 240)
PALE_BLUE = RGBColor(239, 247, 255)
PALE_GREEN = RGBColor(236, 248, 243)
PALE_ORANGE = RGBColor(255, 244, 231)
WHITE = RGBColor(255, 255, 255)


def remove_all_shapes(slide):
    for shape in list(slide.shapes):
        element = shape._element
        element.getparent().remove(element)


def set_background(slide):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = WHITE


def add_text(
    slide,
    text,
    x,
    y,
    w,
    h,
    size=12,
    color=INK,
    bold=False,
    align=PP_ALIGN.LEFT,
    valign=MSO_ANCHOR.MIDDLE,
    margin=0.03,
):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    frame = box.text_frame
    frame.clear()
    frame.word_wrap = True
    frame.margin_left = frame.margin_right = Inches(margin)
    frame.margin_top = frame.margin_bottom = Inches(margin)
    frame.vertical_anchor = valign
    paragraph = frame.paragraphs[0]
    paragraph.text = text
    paragraph.alignment = align
    paragraph.space_after = Pt(0)
    paragraph.line_spacing = 1.05
    run = paragraph.runs[0]
    run.font.name = FONT
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    return box


def add_rect(slide, x, y, w, h, fill, line=None, rounded=False):
    shape_type = MSO_SHAPE.ROUNDED_RECTANGLE if rounded else MSO_SHAPE.RECTANGLE
    shape = slide.shapes.add_shape(
        shape_type, Inches(x), Inches(y), Inches(w), Inches(h)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.color.rgb = line if line else fill
    shape.line.width = Pt(0.7 if line else 0.1)
    if rounded:
        shape.adjustments[0] = 0.08
    return shape


def add_circle(slide, x, y, d, fill):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.OVAL, Inches(x), Inches(y), Inches(d), Inches(d)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.color.rgb = fill
    return shape


def add_line(slide, x1, y1, x2, y2, color=LIGHT, width=1):
    connector = slide.shapes.add_connector(
        MSO_CONNECTOR.STRAIGHT,
        Inches(x1),
        Inches(y1),
        Inches(x2),
        Inches(y2),
    )
    connector.line.color.rgb = color
    connector.line.width = Pt(width)
    return connector


def add_step_card(slide, x, y, w, h, number, title, body, accent):
    add_rect(slide, x, y, w, h, WHITE, LIGHT, rounded=True)
    add_rect(slide, x, y, 0.07, h, accent)
    add_circle(slide, x + 0.18, y + 0.18, 0.54, accent)
    add_text(
        slide,
        number,
        x + 0.18,
        y + 0.29,
        0.54,
        0.19,
        10,
        WHITE,
        True,
        PP_ALIGN.CENTER,
    )
    add_text(slide, title, x + 0.84, y + 0.17, w - 1.05, 0.34, 13.3, INK, True)
    add_text(
        slide,
        body,
        x + 0.23,
        y + 0.72,
        w - 0.46,
        h - 0.84,
        8.35,
        MID,
        False,
        valign=MSO_ANCHOR.TOP,
    )


def set_notes(slide, text):
    notes_slide = slide.notes_slide
    for shape in notes_slide.shapes:
        if not getattr(shape, "is_placeholder", False):
            continue
        try:
            if shape.placeholder_format.type == 2:
                shape.text_frame.clear()
                shape.text_frame.paragraphs[0].text = text
                return
        except (ValueError, AttributeError):
            continue


def build_ppt():
    presentation = Presentation(SOURCE)
    if len(presentation.slides) < 9:
        raise RuntimeError("当前 PPT 不包含第 9 页。")
    slide = presentation.slides[8]
    logo_blob = presentation.slides[7].shapes[1].image.blob

    remove_all_shapes(slide)
    set_background(slide)

    add_rect(slide, 0.32, 0.39, 0.09, 0.28, RED)
    add_rect(slide, 0.32, 0.67, 0.09, 0.20, BLUE)
    add_text(slide, "反思：从异常现象到规则修正", 0.56, 0.30, 7.60, 0.50, 23, INK, True)
    add_text(
        slide,
        "真实案例：支路350（源道路 419，D6 单侧人行道）在 J0109 快速路2合流节点深入路口核心",
        0.58,
        0.86,
        10.70,
        0.32,
        10.5,
        MID,
    )
    slide.shapes.add_picture(
        BytesIO(logo_blob), Inches(10.08), Inches(0.43), width=Inches(2.84)
    )

    steps = [
        (
            "1",
            "发现",
            "全量城市模型在 Blender 中检查时，发现支路350的人行道穿入快速路合流区。小范围样例未覆盖这种“支路 D6 × 快速路 A1”组合。",
            BLUE,
        ),
        (
            "2",
            "隔离",
            "由模型对象追溯 road_name 和 source_road_id=419；在路口语义中定位 J0109，并单独导出 J0109 OBJ/FBX，排除全城模型干扰。",
            CYAN,
        ),
        (
            "3",
            "解释",
            "D6 断面为 5m 车行道 + 单侧 2m 人行道。原逻辑只判断“是否为 D6”，未判断路口是否含快速路，错误保留了支路人行道保护段。",
            ORANGE,
        ),
        (
            "4",
            "修正",
            "增加路口上下文判断：存在 expressway arm 时，非快速路 arm 不参与 one-sided sidewalk protection；支路构件按共享路口面裁剪。",
            RED,
        ),
        (
            "5",
            "检查",
            "先复查 J0109，再回归全量模型；同时对比 J0106，确保普通交叉口仍保留 D6 人行道。共享逻辑使 CIM3/CIM4 同步修复。",
            GREEN,
        ),
    ]
    x_values = [0.58, 3.06, 5.54, 8.02, 10.50]
    for step, x in zip(steps, x_values):
        add_step_card(slide, x, 1.45, 2.25, 3.28, *step)
    for x in [2.84, 5.32, 7.80, 10.28]:
        add_text(slide, "→", x, 2.90, 0.18, 0.27, 15, RGBColor(160, 197, 227), True, PP_ALIGN.CENTER)

    add_text(slide, "规则修正的关键对照", 0.58, 5.04, 1.74, 0.27, 10.5, NAVY, True)
    add_rect(slide, 2.45, 4.93, 4.70, 0.83, PALE_GREEN, PALE_GREEN, rounded=True)
    add_text(slide, "J0106  普通交叉口", 2.68, 5.08, 1.72, 0.25, 10.2, GREEN, True)
    add_text(
        slide,
        "B4 / C1 主干路 + D6 支路，无快速路 arm → 保留单侧人行道连接",
        4.38,
        5.05,
        2.52,
        0.31,
        8.1,
        INK,
    )
    add_rect(slide, 7.38, 4.93, 5.35, 0.83, PALE_ORANGE, PALE_ORANGE, rounded=True)
    add_text(slide, "J0109  快速路合流口", 7.62, 5.08, 1.95, 0.25, 10.2, ORANGE, True)
    add_text(
        slide,
        "A1 快速路 + D6 支路，RAMP_MERGE → 抑制支路保护段，避免深入路口核心",
        9.52,
        5.05,
        2.91,
        0.31,
        8.1,
        INK,
    )

    add_rect(slide, 0.58, 6.05, 12.15, 0.62, NAVY, NAVY, rounded=True)
    add_text(
        slide,
        "方法沉淀：对象追溯 → 单路口隔离 → 几何与语义联合解释 → 上下文规则修正 → 单点与全量回归",
        0.84,
        6.18,
        11.62,
        0.29,
        10.5,
        WHITE,
        True,
        PP_ALIGN.CENTER,
    )
    add_line(slide, 0.58, 6.94, 12.74, 6.94, LIGHT, 0.7)
    add_text(slide, "成果与反思", 0.58, 7.00, 1.50, 0.18, 7.5, MID)
    add_text(
        slide,
        "10",
        12.40,
        7.00,
        0.33,
        0.18,
        8,
        RGBColor(170, 187, 205),
        False,
        PP_ALIGN.RIGHT,
    )

    set_notes(
        slide,
        (
            "这一页用一个真实路口问题说明我们的定位过程。全量城市模型生成后，"
            "在 Blender 中发现支路350的人行道进入了快速路2的合流核心区。"
            "这类问题在小范围样例中没有出现，因为样例没有覆盖D6单侧人行道与A1快速路合流的组合。"
            "定位时先从模型对象名称追溯到road_name和源道路编号419，再查询路口语义，"
            "确认问题发生在J0109。该路口由两条快速路arm和一条D6支路arm组成，类型为RAMP_MERGE。"
            "随后只导出J0109的OBJ和FBX，在Blender中独立检查。"
            "代码回查发现，D6单侧人行道保护逻辑原本用于避免道路端点附近的人行道被路口面完全裁掉，"
            "但旧逻辑只判断断面是否为D6，没有判断当前路口是否包含快速路，"
            "因此错误地把支路350的人行道保护段保留到了快速路核心区。"
            "修正时增加路口上下文判断：如果存在expressway arm，非快速路arm不再参与单侧人行道保护，"
            "并继续使用统一路口面进行平面差集裁剪。"
            "验证时不仅检查J0109问题是否消失，还对比支路另一端的J0106。"
            "J0106是普通交叉口，没有快速路arm，因此仍保留正常的D6单侧人行道连接。"
            "最后再回归全量模型。由于CIM3和CIM4共享同一路口处理逻辑，这次规则修正会同时作用于两个等级。"
            "这个过程说明，城市级自动建模的问题解决不能依赖局部补片，而应建立对象追溯、"
            "单路口隔离、语义解释、上下文规则修正和全量回归的闭环。"
        ),
    )

    presentation.save(OUTPUT)


def font(size, bold=False):
    regular = r"C:\Windows\Fonts\msyh.ttc"
    bold_path = r"C:\Windows\Fonts\msyhbd.ttc"
    return ImageFont.truetype(bold_path if bold else regular, size)


def rounded(draw, xy, fill, outline=(212, 227, 240), radius=14, width=2):
    draw.rounded_rectangle(xy, radius=radius, fill=fill, outline=outline, width=width)


def wrap(draw, text, x, y, max_width, font_obj, fill, line_gap=8):
    lines = []
    current = ""
    for char in text:
        candidate = current + char
        if draw.textlength(candidate, font=font_obj) <= max_width:
            current = candidate
        else:
            lines.append(current)
            current = char
    if current:
        lines.append(current)
    line_height = font_obj.size + line_gap
    for index, line_text in enumerate(lines):
        draw.text((x, y + index * line_height), line_text, fill=fill, font=font_obj)


def build_preview():
    image = Image.new("RGB", (1600, 900), "white")
    draw = ImageDraw.Draw(image)
    draw.rectangle((38, 48, 48, 82), fill=(239, 31, 55))
    draw.rectangle((38, 82, 48, 106), fill=(16, 112, 224))
    draw.text((68, 42), "反思：从异常现象到规则修正", fill=(26, 46, 70), font=font(38, True))
    draw.text(
        (70, 107),
        "真实案例：支路350（源道路 419，D6 单侧人行道）在 J0109 快速路2合流节点深入路口核心",
        fill=(86, 108, 134),
        font=font(18),
    )
    steps = [
        ("1", "发现", "全量模型发现支路350人行道穿入快速路合流区；小样例未覆盖该组合。", (16, 112, 224)),
        ("2", "隔离", "对象追溯至源道路419；语义定位J0109；单独导出OBJ/FBX检查。", (24, 166, 207)),
        ("3", "解释", "旧逻辑只判断D6，未识别路口含快速路，错误保留支路保护段。", (242, 143, 35)),
        ("4", "修正", "含expressway arm时，非快速路arm不参与单侧人行道保护。", (239, 31, 55)),
        ("5", "检查", "J0109单点与全量回归；对照J0106仍保留正常D6连接。", (33, 145, 103)),
    ]
    for index, (num, title, body, accent) in enumerate(steps):
        x = 70 + index * 298
        rounded(draw, (x, 175, x + 270, 565), "white")
        draw.rectangle((x, 184, x + 8, 556), fill=accent)
        draw.ellipse((x + 25, 200, x + 82, 257), fill=accent)
        draw.text((x + 53, 228), num, fill="white", font=font(18, True), anchor="mm")
        draw.text((x + 102, 205), title, fill=(26, 46, 70), font=font(24, True))
        wrap(draw, body, x + 27, 295, 220, font(16), (86, 108, 134), 9)
        if index < 4:
            draw.text((x + 280, 355), "→", fill=(160, 197, 227), font=font(25, True))
    draw.text((70, 605), "规则修正的关键对照", fill=(16, 50, 88), font=font(19, True))
    rounded(draw, (300, 590, 850, 680), (236, 248, 243), (236, 248, 243))
    draw.text((330, 612), "J0106  普通交叉口", fill=(33, 145, 103), font=font(18, True))
    draw.text((330, 644), "无快速路 arm → 保留 D6 单侧人行道连接", fill=(26, 46, 70), font=font(14))
    rounded(draw, (880, 590, 1530, 680), (255, 244, 231), (255, 244, 231))
    draw.text((910, 612), "J0109  快速路合流口", fill=(242, 143, 35), font=font(18, True))
    draw.text((910, 644), "A1 快速路 + D6 支路 → 抑制支路保护段", fill=(26, 46, 70), font=font(14))
    rounded(draw, (70, 730, 1530, 800), (16, 50, 88), (16, 50, 88))
    draw.text(
        (800, 765),
        "对象追溯 → 单路口隔离 → 几何与语义联合解释 → 上下文规则修正 → 单点与全量回归",
        fill="white",
        font=font(20, True),
        anchor="mm",
    )
    draw.line((70, 855, 1530, 855), fill=(212, 227, 240), width=2)
    draw.text((70, 866), "成果与反思", fill=(86, 108, 134), font=font(13))
    draw.text((1505, 866), "10", fill=(170, 187, 205), font=font(14), anchor="ra")
    PREVIEW.parent.mkdir(parents=True, exist_ok=True)
    image.save(PREVIEW)


if __name__ == "__main__":
    build_ppt()
    build_preview()
    print(OUTPUT)
    print(PREVIEW)
