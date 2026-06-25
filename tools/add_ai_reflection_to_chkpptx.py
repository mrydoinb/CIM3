from io import BytesIO
from pathlib import Path

from PIL import Image, ImageDraw, ImageFont
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_CONNECTOR, MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt


ROOT = Path(r"C:\Users\22838\Desktop\chk\CIMAgent\cim_road_poc")
SOURCE = Path(r"C:\Users\22838\Desktop\chk\CIMAgent\PPT汇报\chkpptx.pptx")
OUTPUT = ROOT / "chkpptx_AI反思及后续重点工作完善版.pptx"
PREVIEW = ROOT / "output" / "chkpptx_AI反思_第9页预览.png"
NEXT_PREVIEW = ROOT / "output" / "chkpptx_后续重点工作_第10页预览.png"

FONT = "Microsoft YaHei"

NAVY = RGBColor(16, 50, 88)
DEEP_BLUE = RGBColor(0, 82, 156)
BLUE = RGBColor(16, 112, 224)
CYAN = RGBColor(24, 166, 207)
GREEN = RGBColor(33, 145, 103)
ORANGE = RGBColor(229, 139, 48)
RED = RGBColor(235, 35, 58)
INK = RGBColor(30, 49, 72)
MID = RGBColor(91, 111, 135)
LIGHT = RGBColor(215, 228, 240)
PALE_BLUE = RGBColor(239, 247, 255)
PALE_GREEN = RGBColor(236, 248, 243)
WHITE = RGBColor(255, 255, 255)


def remove_all_shapes(slide):
    for shape in list(slide.shapes):
        element = shape._element
        element.getparent().remove(element)


def set_solid_background(slide, color):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_text(
    slide,
    text,
    x,
    y,
    w,
    h,
    size=12,
    color=INK,
    bold=False,
    align=PP_ALIGN.LEFT,
    valign=MSO_ANCHOR.MIDDLE,
    margin=0.03,
):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    frame = box.text_frame
    frame.clear()
    frame.word_wrap = True
    frame.margin_left = frame.margin_right = Inches(margin)
    frame.margin_top = frame.margin_bottom = Inches(margin)
    frame.vertical_anchor = valign
    paragraph = frame.paragraphs[0]
    paragraph.text = text
    paragraph.alignment = align
    paragraph.space_after = Pt(0)
    paragraph.line_spacing = 1.05
    run = paragraph.runs[0]
    run.font.name = FONT
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    return box


def add_rect(slide, x, y, w, h, fill, line=None, rounded=False):
    shape_type = MSO_SHAPE.ROUNDED_RECTANGLE if rounded else MSO_SHAPE.RECTANGLE
    shape = slide.shapes.add_shape(
        shape_type, Inches(x), Inches(y), Inches(w), Inches(h)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.color.rgb = line if line else fill
    shape.line.width = Pt(0.7 if line else 0.1)
    if rounded:
        shape.adjustments[0] = 0.08
    return shape


def add_circle(slide, x, y, d, fill):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.OVAL, Inches(x), Inches(y), Inches(d), Inches(d)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.color.rgb = fill
    return shape


def add_line(slide, x1, y1, x2, y2, color=LIGHT, width=1):
    shape = slide.shapes.add_connector(
        MSO_CONNECTOR.STRAIGHT,
        Inches(x1),
        Inches(y1),
        Inches(x2),
        Inches(y2),
    )
    shape.line.color.rgb = color
    shape.line.width = Pt(width)
    return shape


def add_label(slide, text, x, y, w, fill, color):
    add_rect(slide, x, y, w, 0.28, fill, fill, rounded=True)
    add_text(
        slide,
        text,
        x + 0.04,
        y + 0.01,
        w - 0.08,
        0.24,
        8.2,
        color,
        True,
        PP_ALIGN.CENTER,
    )


def add_reflection_card(slide, x, y, w, h, number, title, reflection, ai_idea, accent):
    add_rect(slide, x, y, w, h, WHITE, LIGHT, rounded=True)
    add_rect(slide, x, y, 0.08, h, accent)
    add_circle(slide, x + 0.24, y + 0.20, 0.48, accent)
    add_text(
        slide,
        number,
        x + 0.24,
        y + 0.29,
        0.48,
        0.18,
        9.5,
        WHITE,
        True,
        PP_ALIGN.CENTER,
    )
    add_text(slide, title, x + 0.88, y + 0.16, w - 1.16, 0.36, 14, INK, True)

    add_label(slide, "现状反思", x + 0.25, y + 0.72, 0.86, PALE_BLUE, BLUE)
    add_text(
        slide,
        reflection,
        x + 1.25,
        y + 0.67,
        w - 1.50,
        0.50,
        8.9,
        MID,
        False,
        valign=MSO_ANCHOR.TOP,
    )

    add_label(slide, "AI 应用设想", x + 0.25, y + 1.30, 0.98, PALE_GREEN, GREEN)
    add_text(
        slide,
        ai_idea,
        x + 1.35,
        y + 1.23,
        w - 1.60,
        h - 1.33,
        8.9,
        INK,
        False,
        valign=MSO_ANCHOR.TOP,
    )


def add_work_card(slide, x, y, w, h, number, title, items, accent):
    add_rect(slide, x, y, w, h, WHITE, LIGHT, rounded=True)
    add_rect(slide, x, y, 0.08, h, accent)
    add_circle(slide, x + 0.25, y + 0.22, 0.50, accent)
    add_text(
        slide,
        number,
        x + 0.25,
        y + 0.31,
        0.50,
        0.18,
        9.5,
        WHITE,
        True,
        PP_ALIGN.CENTER,
    )
    add_text(slide, title, x + 0.92, y + 0.17, w - 1.22, 0.36, 13.5, INK, True)
    for index, item in enumerate(items):
        item_y = y + 0.70 + index * 0.37
        add_circle(slide, x + 0.31, item_y + 0.08, 0.12, accent)
        add_text(
            slide,
            item,
            x + 0.55,
            item_y,
            w - 0.80,
            0.29,
            8.65,
            MID,
            False,
            valign=MSO_ANCHOR.TOP,
        )


def set_notes(slide, text):
    notes_slide = slide.notes_slide
    for shape in notes_slide.shapes:
        if not getattr(shape, "is_placeholder", False):
            continue
        try:
            if shape.placeholder_format.type == 2:  # BODY / notes placeholder
                frame = shape.text_frame
                frame.clear()
                frame.paragraphs[0].text = text
                return
        except (ValueError, AttributeError):
            continue


def build_ppt():
    presentation = Presentation(SOURCE)
    reflection_slide = presentation.slides[8]
    source_logo = presentation.slides[7].shapes[1].image.blob

    remove_all_shapes(reflection_slide)
    set_solid_background(reflection_slide, WHITE)

    add_text(
        reflection_slide,
        "成果反思与 AI 赋能思考",
        0.56,
        0.30,
        6.90,
        0.50,
        23,
        INK,
        True,
    )
    add_text(
        reflection_slide,
        "AI 不替代确定性建模内核，重点增强数据理解、规则构建、质量检查与知识复用",
        0.58,
        0.86,
        9.70,
        0.32,
        10.5,
        MID,
    )
    add_rect(reflection_slide, 0.32, 0.39, 0.09, 0.28, RED)
    add_rect(reflection_slide, 0.32, 0.67, 0.09, 0.20, BLUE)
    reflection_slide.shapes.add_picture(
        BytesIO(source_logo), Inches(10.08), Inches(0.43), width=Inches(2.84)
    )

    cards = [
        (
            "01",
            "多源数据理解与治理",
            "不同来源数据在字段、分类、精度和业务口径上存在差异，规则匹配仍依赖人工判断。",
            "利用语义模型辅助字段映射、对象分类、异常识别和缺失信息补全建议，并保留置信度与人工复核。",
            BLUE,
        ),
        (
            "02",
            "专业规则解析与配置",
            "专业规范、参数表与经验规则较为分散，转化为机器可执行配置的成本较高。",
            "从规范和设计资料中提取构件、参数与约束，形成规则草案和版本差异，由专家审定后进入规则库。",
            CYAN,
        ),
        (
            "03",
            "模型检查与问题定位",
            "路口缝隙、重叠、构件缺失等问题，目前仍主要依靠人工浏览模型进行发现。",
            "融合几何指标、语义属性与渲染图像开展智能巡检，辅助回溯数据或规则来源并提出修正建议。",
            ORANGE,
        ),
        (
            "04",
            "知识沉淀与专题辅助",
            "规则依据、问题案例和处理经验分散在代码、文档及个人经验中，跨项目复用能力不足。",
            "构建“规范—规则—对象—问题案例”知识库，支持自然语言查询、相似案例检索、参数建议和专题分析。",
            GREEN,
        ),
    ]
    positions = [
        (0.58, 1.40),
        (6.75, 1.40),
        (0.58, 3.55),
        (6.75, 3.55),
    ]
    for card, (x, y) in zip(cards, positions):
        add_reflection_card(
            reflection_slide,
            x,
            y,
            5.98,
            1.88,
            card[0],
            card[1],
            card[2],
            card[3],
            card[4],
        )

    add_rect(reflection_slide, 0.58, 5.72, 12.15, 0.78, NAVY, NAVY, rounded=True)
    add_text(
        reflection_slide,
        "AI 应用边界",
        0.88,
        5.91,
        1.15,
        0.30,
        11.5,
        WHITE,
        True,
        PP_ALIGN.CENTER,
    )
    add_text(
        reflection_slide,
        "AI 负责“理解、建议与检查”",
        2.20,
        5.92,
        2.42,
        0.28,
        10.5,
        RGBColor(195, 225, 255),
        True,
        PP_ALIGN.CENTER,
    )
    add_text(
        reflection_slide,
        "规则与几何引擎负责“确定性生成”",
        4.79,
        5.92,
        3.10,
        0.28,
        10.5,
        RGBColor(195, 238, 225),
        True,
        PP_ALIGN.CENTER,
    )
    add_text(
        reflection_slide,
        "专家负责“审定、追溯与发布”",
        8.08,
        5.92,
        2.72,
        0.28,
        10.5,
        RGBColor(255, 225, 190),
        True,
        PP_ALIGN.CENTER,
    )
    add_text(
        reflection_slide,
        "形成可解释、可追溯、人在回路的人机协同机制",
        10.86,
        5.92,
        1.57,
        0.30,
        8.0,
        WHITE,
        True,
        PP_ALIGN.CENTER,
    )
    add_line(reflection_slide, 0.58, 6.86, 12.74, 6.86, LIGHT, 0.7)
    add_text(
        reflection_slide,
        "成果与反思",
        0.58,
        6.94,
        1.60,
        0.20,
        7.5,
        MID,
    )
    add_text(
        reflection_slide,
        "09",
        12.40,
        6.94,
        0.33,
        0.20,
        8,
        RGBColor(170, 187, 205),
        False,
        PP_ALIGN.RIGHT,
    )

    set_notes(
        reflection_slide,
        (
            "从当前成果看，规则化建模的基本链路已经能够贯通，但也暴露出四方面问题。"
            "第一，多源数据的字段、分类和精度仍需要大量人工理解；第二，专业规范和经验规则转化为配置的成本较高；"
            "第三，模型质量问题主要依靠人工浏览发现；第四，规则依据和问题处理经验尚未形成系统化知识沉淀。"
            "下一步可以考虑引入人工智能，但需要明确边界。AI更适合承担规范理解、数据映射、异常发现、"
            "规则草案生成和知识检索等辅助工作；确定性的几何生成、对象编码以及正式成果发布，"
            "仍应由规则引擎和专家审核控制。最终形成可解释、可追溯、人在回路的人机协同机制。"
        ),
    )

    next_slide = presentation.slides[9]
    remove_all_shapes(next_slide)
    set_solid_background(next_slide, WHITE)

    add_text(
        next_slide,
        "后续重点工作",
        0.56,
        0.30,
        4.20,
        0.50,
        23,
        INK,
        True,
    )
    add_text(
        next_slide,
        "围绕“对象生产—数据沉淀—智能增强—专题应用”形成可持续演进的技术能力体系",
        0.58,
        0.86,
        10.10,
        0.32,
        10.5,
        MID,
    )
    add_rect(next_slide, 0.32, 0.39, 0.09, 0.28, RED)
    add_rect(next_slide, 0.32, 0.67, 0.09, 0.20, BLUE)
    next_slide.shapes.add_picture(
        BytesIO(source_logo), Inches(10.08), Inches(0.43), width=Inches(2.84)
    )

    work_cards = [
        (
            "01",
            "共性自动化建模能力",
            [
                "完善公共道路、区间隧道既有链路，深化路口与专业构件表达",
                "拓展地下管线、公交站、地铁站点等对象的规则化生成",
                "统一对象编码、CIM3/CIM4 分级、专业属性与空间关系",
                "推动断面、构件和语义规则配置化、模板化、版本化",
            ],
            BLUE,
        ),
        (
            "02",
            "超融合数据库研发",
            [
                "统一组织 GIS、CAD、三维模型、业务表、文档和专题数据",
                "研究跨专业统一对象模型、分类编码和多层级表达",
                "建立空间、时间、属性、关系与语义标签的复合索引",
                "完善对象关系、历史版本、增量更新和标准数据服务",
            ],
            CYAN,
        ),
        (
            "03",
            "AI 能力增强",
            [
                "开展规范解析、字段映射、对象分类和规则草案生成",
                "研究融合几何、语义与渲染图像的模型智能检查",
                "建设规范—规则—对象—问题案例的专业知识库",
                "坚持人在回路，保留规则依据、置信度和人工审定过程",
            ],
            GREEN,
        ),
        (
            "04",
            "数字孪生专项建设",
            [
                "构建道路、轨道、管线、站点和地下空间专题场景",
                "建立场景对象与数据库对象、业务属性和专题数据关联",
                "形成专题图层、指标口径、查询统计与关系分析能力",
                "沉淀可复用的场景模板、对象接口和专题应用方法",
            ],
            ORANGE,
        ),
    ]
    work_positions = [
        (0.58, 1.40),
        (6.75, 1.40),
        (0.58, 3.46),
        (6.75, 3.46),
    ]
    for card, (x, y) in zip(work_cards, work_positions):
        add_work_card(
            next_slide,
            x,
            y,
            5.98,
            1.78,
            card[0],
            card[1],
            card[2],
            card[3],
        )

    add_text(
        next_slide,
        "推进次序",
        0.58,
        5.52,
        0.90,
        0.27,
        10.5,
        INK,
        True,
    )
    phases = [
        (
            "近期",
            "统一标准、对象编码与规则模板，补齐道路、隧道和管线关键能力",
            PALE_BLUE,
            BLUE,
            1.60,
            3.35,
        ),
        (
            "中期",
            "推进超融合数据库对象、索引和版本模型，形成 AI 辅助工具链",
            PALE_GREEN,
            GREEN,
            5.10,
            3.38,
        ),
        (
            "持续",
            "以数字孪生专项牵引数据、建模、AI 和专题应用协同演进",
            RGBColor(255, 244, 231),
            ORANGE,
            8.64,
            4.09,
        ),
    ]
    for label, body, fill, color, x, width in phases:
        add_rect(next_slide, x, 5.43, width, 0.66, fill, fill, rounded=True)
        add_text(
            next_slide,
            label,
            x + 0.13,
            5.53,
            0.52,
            0.25,
            9,
            color,
            True,
            PP_ALIGN.CENTER,
        )
        add_text(
            next_slide,
            body,
            x + 0.73,
            5.49,
            width - 0.87,
            0.34,
            7.9,
            INK,
            False,
            valign=MSO_ANCHOR.TOP,
        )

    add_rect(next_slide, 0.58, 6.27, 12.15, 0.48, NAVY, NAVY, rounded=True)
    add_text(
        next_slide,
        "以自动化建模形成数字对象，以超融合数据库沉淀数据资产，以 AI 增强理解与检查，以数字孪生专项牵引应用。",
        0.82,
        6.36,
        11.67,
        0.27,
        10.5,
        WHITE,
        True,
        PP_ALIGN.CENTER,
    )
    add_line(next_slide, 0.58, 6.94, 12.74, 6.94, LIGHT, 0.7)
    add_text(next_slide, "后续重点工作", 0.58, 7.00, 1.60, 0.18, 7.5, MID)
    add_text(
        next_slide,
        "10",
        12.40,
        7.00,
        0.33,
        0.18,
        8,
        RGBColor(170, 187, 205),
        False,
        PP_ALIGN.RIGHT,
    )

    set_notes(
        next_slide,
        (
            "后续重点工作围绕四条主线推进。第一，继续完善共性自动化建模能力，"
            "既要深化公共道路和区间隧道，也要逐步拓展地下管线、公交站和地铁站点，"
            "并统一对象编码、CIM分级、专业属性和空间关系。第二，推进超融合数据库研发，"
            "重点研究统一对象模型、时空复合索引、关系与版本模型以及标准数据服务。"
            "第三，引入人工智能增强规范解析、数据治理、模型检查和知识复用，但保持人在回路，"
            "不替代确定性规则与几何生成。第四，以数字孪生专项牵引场景、数据、指标和专题应用建设。"
            "推进次序上，近期先稳定标准和规则，中期形成数据库与AI辅助工具链，"
            "持续通过数字孪生专项促进各项能力协同演进。"
        ),
    )

    presentation.save(OUTPUT)


def build_preview():
    width, height = 1600, 900
    image = Image.new("RGB", (width, height), "white")
    draw = ImageDraw.Draw(image)
    regular = r"C:\Windows\Fonts\msyh.ttc"
    bold = r"C:\Windows\Fonts\msyhbd.ttc"

    def font(size, is_bold=False):
        return ImageFont.truetype(bold if is_bold else regular, size)

    def rounded_box(xy, fill, outline, radius=14, width_px=2):
        draw.rounded_rectangle(xy, radius=radius, fill=fill, outline=outline, width=width_px)

    draw.rectangle((38, 48, 48, 82), fill=(235, 35, 58))
    draw.rectangle((38, 82, 48, 106), fill=(16, 112, 224))
    draw.text((68, 42), "成果反思与 AI 赋能思考", fill=(30, 49, 72), font=font(39, True))
    draw.text(
        (70, 107),
        "AI 不替代确定性建模内核，重点增强数据理解、规则构建、质量检查与知识复用",
        fill=(91, 111, 135),
        font=font(20),
    )

    cards = [
        ("01", "多源数据理解与治理", "多源数据口径不一，规则匹配依赖人工判断。",
         "AI 辅助字段映射、对象分类、异常识别与缺失信息补全建议。", (16, 112, 224)),
        ("02", "专业规则解析与配置", "规范、参数表和经验规则分散，配置成本较高。",
         "AI 从规范资料提取构件、参数与约束，形成待审定的规则草案。", (24, 166, 207)),
        ("03", "模型检查与问题定位", "缝隙、重叠和构件缺失仍主要依靠人工浏览发现。",
         "融合几何、语义与渲染图开展智能巡检，辅助定位问题来源。", (229, 139, 48)),
        ("04", "知识沉淀与专题辅助", "规则依据和问题经验分散，跨项目复用能力不足。",
         "建设规范—规则—对象—案例知识库，支持检索、建议与专题分析。", (33, 145, 103)),
    ]
    coords = [(70, 170), (815, 170), (70, 440), (815, 440)]
    for (num, title, reflection, idea, accent), (x, y) in zip(cards, coords):
        rounded_box((x, y, x + 690, y + 220), (255, 255, 255), (215, 228, 240))
        draw.rectangle((x, y + 8, x + 9, y + 212), fill=accent)
        draw.ellipse((x + 28, y + 26, x + 82, y + 80), fill=accent)
        draw.text((x + 42, y + 39), num, fill="white", font=font(16, True), anchor="mm")
        draw.text((x + 105, y + 24), title, fill=(30, 49, 72), font=font(25, True))
        draw.rounded_rectangle((x + 28, y + 92, x + 130, y + 122), radius=10, fill=(239, 247, 255))
        draw.text((x + 79, y + 107), "现状反思", fill=(16, 112, 224), font=font(15, True), anchor="mm")
        draw.text((x + 150, y + 92), reflection, fill=(91, 111, 135), font=font(16))
        draw.rounded_rectangle((x + 28, y + 148, x + 145, y + 180), radius=10, fill=(236, 248, 243))
        draw.text((x + 86, y + 164), "AI 应用设想", fill=(33, 145, 103), font=font(15, True), anchor="mm")
        draw.text((x + 165, y + 145), idea, fill=(30, 49, 72), font=font(16))

    rounded_box((70, 705, 1530, 810), (16, 50, 88), (16, 50, 88), radius=16)
    draw.text((105, 742), "AI 应用边界", fill="white", font=font(21, True))
    draw.text((310, 742), "AI：理解、建议与检查", fill=(195, 225, 255), font=font(20, True))
    draw.text((670, 742), "规则与几何引擎：确定性生成", fill=(195, 238, 225), font=font(20, True))
    draw.text((1110, 742), "专家：审定、追溯与发布", fill=(255, 225, 190), font=font(20, True))
    draw.line((70, 855, 1530, 855), fill=(215, 228, 240), width=2)
    draw.text((70, 866), "成果与反思", fill=(91, 111, 135), font=font(13))
    draw.text((1505, 866), "09", fill=(170, 187, 205), font=font(14), anchor="ra")
    PREVIEW.parent.mkdir(parents=True, exist_ok=True)
    image.save(PREVIEW)

    next_image = Image.new("RGB", (width, height), "white")
    draw = ImageDraw.Draw(next_image)
    draw.rectangle((38, 48, 48, 82), fill=(235, 35, 58))
    draw.rectangle((38, 82, 48, 106), fill=(16, 112, 224))
    draw.text((68, 42), "后续重点工作", fill=(30, 49, 72), font=font(39, True))
    draw.text(
        (70, 107),
        "围绕“对象生产—数据沉淀—智能增强—专题应用”形成可持续演进的技术能力体系",
        fill=(91, 111, 135),
        font=font(20),
    )

    preview_cards = [
        ("01", "共性自动化建模能力", [
            "深化道路、隧道既有链路与关键构件表达",
            "拓展管线、公交站、地铁站等对象生成",
            "统一对象编码、CIM 分级与空间关系",
            "规则配置化、模板化与版本化",
        ], (16, 112, 224)),
        ("02", "超融合数据库研发", [
            "统一组织 GIS、CAD、模型、表格和文档",
            "研究统一对象模型与多层级表达",
            "建立时空、属性、关系和语义复合索引",
            "完善版本、增量更新和数据服务",
        ], (24, 166, 207)),
        ("03", "AI 能力增强", [
            "规范解析、字段映射与规则草案生成",
            "融合几何、语义和图像开展智能检查",
            "建设规范—规则—对象—案例知识库",
            "坚持人在回路与专家审定",
        ], (33, 145, 103)),
        ("04", "数字孪生专项建设", [
            "建设多类基础设施专题场景",
            "关联场景对象、数据库对象与业务数据",
            "形成专题图层、指标和关系分析能力",
            "沉淀场景模板、对象接口和应用方法",
        ], (229, 139, 48)),
    ]
    coords = [(70, 170), (815, 170), (70, 420), (815, 420)]
    for (num, title, items, accent), (x, y) in zip(preview_cards, coords):
        rounded_box((x, y, x + 690, y + 210), (255, 255, 255), (215, 228, 240))
        draw.rectangle((x, y + 8, x + 9, y + 202), fill=accent)
        draw.ellipse((x + 28, y + 24, x + 82, y + 78), fill=accent)
        draw.text((x + 55, y + 51), num, fill="white", font=font(16, True), anchor="mm")
        draw.text((x + 105, y + 23), title, fill=(30, 49, 72), font=font(24, True))
        for index, item in enumerate(items):
            item_y = y + 83 + index * 30
            draw.ellipse((x + 34, item_y + 7, x + 44, item_y + 17), fill=accent)
            draw.text((x + 58, item_y), item, fill=(91, 111, 135), font=font(15))

    draw.text((70, 654), "推进次序", fill=(30, 49, 72), font=font(19, True))
    preview_phases = [
        ("近期", "统一标准、编码和规则模板，补齐关键对象能力", (239, 247, 255), (16, 112, 224), 160, 430),
        ("中期", "推进数据库对象、索引、版本模型与 AI 工具链", (236, 248, 243), (33, 145, 103), 610, 430),
        ("持续", "以数字孪生专项牵引各项能力协同演进", (255, 244, 231), (229, 139, 48), 1060, 470),
    ]
    for label, body, fill, color, x, box_w in preview_phases:
        draw.rounded_rectangle((x, 640, x + box_w, 710), radius=12, fill=fill)
        draw.text((x + 20, 660), label, fill=color, font=font(17, True))
        draw.text((x + 85, 657), body, fill=(30, 49, 72), font=font(14))

    draw.rounded_rectangle((70, 745, 1530, 810), radius=14, fill=(16, 50, 88))
    draw.text(
        (800, 777),
        "自动化建模形成对象 · 超融合数据库沉淀资产 · AI 增强理解检查 · 数字孪生牵引应用",
        fill="white",
        font=font(20, True),
        anchor="mm",
    )
    draw.line((70, 855, 1530, 855), fill=(215, 228, 240), width=2)
    draw.text((70, 866), "后续重点工作", fill=(91, 111, 135), font=font(13))
    draw.text((1505, 866), "10", fill=(170, 187, 205), font=font(14), anchor="ra")
    next_image.save(NEXT_PREVIEW)


if __name__ == "__main__":
    build_ppt()
    build_preview()
    print(OUTPUT)
    print(PREVIEW)
    print(NEXT_PREVIEW)
